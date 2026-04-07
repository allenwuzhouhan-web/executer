#!/usr/bin/env python3
"""
Executer PPT Engine
Creates .pptx files from JSON spec + optional design language.

Usage:
    python3 ppt_engine.py --spec spec.json --output output.pptx
    python3 ppt_engine.py --spec spec.json --design design_language.json --output output.pptx
    echo '{"slides":[...]}' | python3 ppt_engine.py --stdin --output output.pptx
"""

import argparse
import importlib.util
import json
import os
import re
import sys
import tempfile
from pathlib import Path

# Import image_utils from same directory
try:
    from image_utils import resolve_image, cleanup_temp_images
except ImportError:
    _iu_path = Path(__file__).parent / "image_utils.py"
    if _iu_path.exists():
        _spec = importlib.util.spec_from_file_location("image_utils", _iu_path)
        _mod = importlib.util.module_from_spec(_spec)
        _spec.loader.exec_module(_mod)
        resolve_image = _mod.resolve_image
        cleanup_temp_images = _mod.cleanup_temp_images
    else:
        resolve_image = lambda source, temp_dir=None: source if source and Path(source).exists() else None
        cleanup_temp_images = lambda d: None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
except ImportError:
    print(json.dumps({"success": False, "error": "python-pptx not installed. Run: pip3 install python-pptx"}))
    sys.exit(1)


# ─── Defaults ───────────────────────────────────────────────────────────────

DEFAULTS = {
    "font_primary": "Helvetica Neue",
    "font_secondary": "Helvetica Neue",
    "title_size": 40,
    "subtitle_size": 24,
    "heading_size": 28,
    "body_size": 18,
    "bullet_size": 16,
    "caption_size": 12,
    "note_size": 10,
    "color_text": "1E293B",               # Slate-800
    "color_text_secondary": "64748B",      # Slate-500
    "color_body": "334155",                # Slate-700
    "color_accent": "2563EB",              # Blue-600
    "color_accent2": "0891B2",             # Cyan-600
    "color_bg": "FFFFFF",
    "color_light_gray": "F1F5F9",          # Slate-100
    "color_mid_gray": "CBD5E1",            # Slate-300
    "color_slide_number": "94A3B8",        # Slate-400
    "slide_width": Inches(13.333),
    "slide_height": Inches(7.5),
    "margin_left_pct": 5,
    "margin_right_pct": 5,
    "margin_top_pct": 12,
    "content_top_pct": 22,
    "line_spacing": 1.2,
    "color_accent_light": "BFDBFE",        # Blue-200
    "color_accent_dark": "1E40AF",         # Blue-800
    "color_bg_subtle": "EFF6FF",           # Blue-50
    "color_accent2_light": "A5F3FC",       # Cyan-200
}


def hex_to_rgb(hex_str):
    """Convert '#AABBCC' or 'AABBCC' to RGBColor."""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        return RGBColor(0x1A, 0x1A, 0x1A)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _parse_hex(hex_str):
    """Parse hex color to (r, g, b) ints."""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        return (0x1A, 0x1A, 0x1A)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def tint_color(hex_str, factor=0.7):
    """Blend color toward white. factor=0.7 → 70% toward white (very light)."""
    r, g, b = _parse_hex(hex_str)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"{r:02X}{g:02X}{b:02X}"


def shade_color(hex_str, factor=0.3):
    """Blend color toward black. factor=0.3 → 30% darker."""
    r, g, b = _parse_hex(hex_str)
    r = int(r * (1 - factor))
    g = int(g * (1 - factor))
    b = int(b * (1 - factor))
    return f"{r:02X}{g:02X}{b:02X}"


def _srgb_to_linear(c):
    """Convert sRGB channel (0-255) to linear for WCAG luminance."""
    c = c / 255.0
    return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4


def wcag_relative_luminance(hex_str):
    """Calculate WCAG 2.0 relative luminance from hex color."""
    r, g, b = _parse_hex(hex_str)
    return 0.2126 * _srgb_to_linear(r) + 0.7152 * _srgb_to_linear(g) + 0.0722 * _srgb_to_linear(b)


def wcag_contrast_ratio(fg_hex, bg_hex):
    """Calculate WCAG 2.0 contrast ratio between two hex colors."""
    l1 = wcag_relative_luminance(fg_hex)
    l2 = wcag_relative_luminance(bg_hex)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def ensure_contrast(fg_hex, bg_hex, min_ratio=4.5):
    """Adjust fg color to meet WCAG contrast ratio against bg.

    For normal text (< 18pt): min_ratio=4.5
    For large text (>= 18pt bold or >= 24pt): min_ratio=3.0
    Returns the adjusted foreground hex (without '#').
    """
    fg_hex = fg_hex.lstrip("#")
    bg_hex = bg_hex.lstrip("#")

    ratio = wcag_contrast_ratio(fg_hex, bg_hex)
    if ratio >= min_ratio:
        return fg_hex  # Already passes

    # Determine direction: if bg is light, darken fg; if bg is dark, lighten fg
    bg_lum = wcag_relative_luminance(bg_hex)
    r, g, b = _parse_hex(fg_hex)

    # Try up to 25 steps of adjustment
    for step in range(1, 26):
        factor = step * 0.04  # 4% per step, up to 100%
        if bg_lum > 0.5:
            # Dark text on light bg — darken the text
            nr = int(r * (1 - factor))
            ng = int(g * (1 - factor))
            nb = int(b * (1 - factor))
        else:
            # Light text on dark bg — lighten the text
            nr = int(r + (255 - r) * factor)
            ng = int(g + (255 - g) * factor)
            nb = int(b + (255 - b) * factor)
        nr = max(0, min(255, nr))
        ng = max(0, min(255, ng))
        nb = max(0, min(255, nb))
        candidate = f"{nr:02X}{ng:02X}{nb:02X}"
        if wcag_contrast_ratio(candidate, bg_hex) >= min_ratio:
            return candidate

    # Fallback: pure black or white
    return "000000" if bg_lum > 0.5 else "FFFFFF"


def pct_to_emu(pct, total_emu):
    """Convert percentage to EMU based on total dimension."""
    return int(total_emu * pct / 100)


def log(msg):
    """Log to stderr (stdout is reserved for JSON result)."""
    print(msg, file=sys.stderr)


# ─── Design Language Loader ────────────────────────────────────────────────

class DesignLanguage:
    """Loads and provides design tokens from extracted design_language.json."""

    def __init__(self, path=None):
        self.data = None
        self.d = dict(DEFAULTS)
        self._layout_patterns = {}
        if path and Path(path).exists():
            try:
                self.data = json.loads(Path(path).read_text(encoding="utf-8"))
                self._apply()
                log(f"Loaded design language from: {path}")
            except Exception as e:
                log(f"Warning: Failed to parse design language: {e}")

    def _valid_hex(self, h):
        """Validate that a string is a valid 6-digit hex color."""
        h = h.lstrip("#")
        return bool(re.match(r'^[0-9A-Fa-f]{6}$', h))

    def _apply(self):
        ds = self.data.get("design_system", {})

        # === SEMANTIC COLORS (preferred — context-aware classification) ===
        semantic = ds.get("semantic_colors", {})
        semantic_applied = False
        if semantic:
            if semantic.get("text_primary") and self._valid_hex(semantic["text_primary"]):
                self.d["color_text"] = semantic["text_primary"].lstrip("#")
                self.d["color_body"] = semantic["text_primary"].lstrip("#")
                semantic_applied = True
            if semantic.get("text_secondary") and self._valid_hex(semantic["text_secondary"]):
                self.d["color_text_secondary"] = semantic["text_secondary"].lstrip("#")
            if semantic.get("background") and self._valid_hex(semantic["background"]):
                self.d["color_bg"] = semantic["background"].lstrip("#")
                # Derive light_gray and mid_gray from background luminance
                bg_lum = self._luminance(semantic["background"])
                if bg_lum > 200:  # Light background
                    self.d["color_light_gray"] = "F5F5F5"
                    self.d["color_mid_gray"] = "DDDDDD"
                else:  # Dark background
                    self.d["color_light_gray"] = "2A2A2A"
                    self.d["color_mid_gray"] = "444444"
                semantic_applied = True
            if semantic.get("accent") and self._valid_hex(semantic["accent"]):
                self.d["color_accent"] = semantic["accent"].lstrip("#")
                semantic_applied = True
            if semantic.get("accent2") and self._valid_hex(semantic["accent2"]):
                self.d["color_accent2"] = semantic["accent2"].lstrip("#")

        # === FALLBACK: Raw color palette (legacy, less accurate) ===
        if not semantic_applied:
            palette = ds.get("color_palette", [])
            if len(palette) >= 1 and self._valid_hex(palette[0].get("hex", "")):
                self.d["color_accent"] = palette[0]["hex"].lstrip("#")
            if len(palette) >= 2 and self._valid_hex(palette[1].get("hex", "")):
                self.d["color_accent2"] = palette[1]["hex"].lstrip("#")
            if len(palette) >= 3 and self._valid_hex(palette[2].get("hex", "")):
                self.d["color_text"] = palette[2]["hex"].lstrip("#") if self._is_dark(palette[2]["hex"]) else self.d["color_text"]

        # Theme colors as FALLBACK only — semantic colors take priority
        if not semantic_applied:
            theme = self.data.get("theme", {}).get("color_scheme", {})
            if theme.get("dk1") and self._valid_hex(theme["dk1"]):
                self.d["color_text"] = theme["dk1"].lstrip("#")
            if theme.get("lt1") and self._valid_hex(theme["lt1"]):
                self.d["color_bg"] = theme["lt1"].lstrip("#")
            if theme.get("accent1") and self._valid_hex(theme["accent1"]):
                self.d["color_accent"] = theme["accent1"].lstrip("#")
            if theme.get("accent2") and self._valid_hex(theme["accent2"]):
                self.d["color_accent2"] = theme["accent2"].lstrip("#")

        # Typography
        typo = ds.get("typography", {})
        fonts = typo.get("fonts_by_frequency", [])
        if len(fonts) >= 1:
            self.d["font_primary"] = fonts[0]["font"]
        if len(fonts) >= 2:
            self.d["font_secondary"] = fonts[1]["font"]
        else:
            self.d["font_secondary"] = self.d["font_primary"]

        # Text hierarchy — clamp to sane ranges so decorative giant text
        # (e.g., 310pt "HISTORY") doesn't break layouts. The engine's layout
        # builders are designed for normal presentation sizes.
        hierarchy = ds.get("text_hierarchy", [])
        role_map = {
            "title": ("title_size", 24, 56),
            "subtitle": ("subtitle_size", 16, 36),
            "heading": ("heading_size", 18, 40),
            "subheading": ("heading_size", 14, 32),
            "body": ("body_size", 12, 24),
            "caption": ("caption_size", 8, 16),
            "footnote": ("note_size", 6, 14),
        }
        for item in hierarchy:
            role = item.get("likely_role", "")
            if role in role_map and item.get("size_pt"):
                key, min_pt, max_pt = role_map[role]
                size = item["size_pt"]
                # Decorative text (>80pt) is artistic, not functional — scale down
                if size > 80:
                    size = max_pt  # Use max of the sane range
                self.d[key] = max(min_pt, min(max_pt, size))

        # Slide dimensions
        dim = self.data.get("slide_dimensions", {})
        if dim.get("width_emu") and dim.get("height_emu"):
            self.d["slide_width"] = dim["width_emu"]
            self.d["slide_height"] = dim["height_emu"]

        # Global spacing from extracted layout patterns
        spacing = self.data.get("global_spacing", {})
        if spacing.get("margin_left_pct"):
            self.d["margin_left_pct"] = spacing["margin_left_pct"]
        if spacing.get("margin_right_pct"):
            self.d["margin_right_pct"] = spacing["margin_right_pct"]
        if spacing.get("margin_top_pct"):
            self.d["margin_top_pct"] = spacing["margin_top_pct"]
        if spacing.get("content_top_pct"):
            self.d["content_top_pct"] = spacing["content_top_pct"]

        # Store layout patterns for per-layout position overrides
        self._layout_patterns = self.data.get("layout_patterns", {})

        # Visual effects — ONLY enabled when the source deck actually uses them.
        # Default is OFF for everything. Clean, minimal design by default.
        # Now captures PARAMETERS so effects match the source deck's style, not generic defaults.
        ve = self.data.get("visual_effects", {}) if self.data else {}
        self._use_shadows = ve.get("has_shadows", False)
        self._use_gradients = ve.get("has_gradients", False)
        self._use_rounded = ve.get("has_rounded_corners", False) or \
                            "ROUNDED_RECTANGLE" in str(ve.get("shape_variety", []))
        self._decorative_elements = ve.get("decorative_elements", [])
        # Shadow parameters from source deck (used by _add_shadow)
        self._shadow_style = ve.get("shadow_style", {})
        # Corner radius from source deck
        self._corner_radius_pct = ve.get("corner_radius_pct")
        # Gradient angle from source deck
        self._gradient_angles = ve.get("gradient_angles", [])

        # Tint/shade system — derived automatically from accent colors
        accent = self.d.get("color_accent", "444444")
        accent2 = self.d.get("color_accent2", accent)
        self.d["color_accent_light"] = tint_color(accent, 0.75)
        self.d["color_accent_dark"] = shade_color(accent, 0.3)
        self.d["color_bg_subtle"] = tint_color(accent, 0.92)
        self.d["color_accent2_light"] = tint_color(accent2, 0.75)

        # Design philosophy — informs layout behavior
        dp = self.data.get("design_philosophy", {}) if self.data else {}
        self._content_density = dp.get("content_density", "moderate")
        self._whitespace_style = dp.get("whitespace_style", "balanced")
        self._layout_complexity = dp.get("layout_complexity", "clean")
        self._dominant_alignment = dp.get("dominant_alignment", "left")
        self._color_restraint = dp.get("color_restraint", "restrained")
        self._effects_usage = dp.get("effects_usage", "minimal")

        # Apply philosophy to rendering parameters
        self._apply_philosophy()

    def _apply_philosophy(self):
        """Adjust rendering parameters based on design philosophy values."""
        # --- content_density: affects max bullets and font sizes ---
        if self._content_density == "sparse":
            # Fewer items, larger text for breathing room
            self._max_bullets = 4
            self._font_size_adjust = 2  # bump all body/bullet sizes up
        elif self._content_density == "dense":
            # Allow more items, slightly smaller text
            self._max_bullets = 8
            self._font_size_adjust = -1
        else:  # moderate
            self._max_bullets = 6
            self._font_size_adjust = 0

        # Apply font size adjustment to body/bullet sizes
        if self._font_size_adjust:
            for key in ("body_size", "bullet_size"):
                self.d[key] = max(10, self.d[key] + self._font_size_adjust)

        # --- whitespace_style: affects margins and padding ---
        if self._whitespace_style == "generous":
            self._margin_adjust = 3   # add 3% to each margin
            self._card_padding_factor = 1.4
        elif self._whitespace_style == "tight":
            self._margin_adjust = -2  # reduce margins
            self._card_padding_factor = 0.7
        else:  # balanced
            self._margin_adjust = 0
            self._card_padding_factor = 1.0

        if self._margin_adjust:
            for key in ("margin_left_pct", "margin_right_pct"):
                self.d[key] = max(2, self.d[key] + self._margin_adjust)
            self.d["margin_top_pct"] = max(4, self.d["margin_top_pct"] + self._margin_adjust)
            self.d["content_top_pct"] = max(10, self.d["content_top_pct"] + self._margin_adjust)

        # --- effects_usage: override shadow/gradient flags ---
        if self._effects_usage == "minimal":
            self._use_shadows = False
            self._use_gradients = False
        elif self._effects_usage == "heavy":
            self._use_shadows = True
            # gradients stay as detected — heavy just ensures shadows are on

        # --- color_restraint: limit accent color usage ---
        if self._color_restraint == "monochrome":
            # Force accent2 to match accent (single-color palette)
            self.d["color_accent2"] = self.d["color_accent"]
            self.d["color_accent2_light"] = self.d["color_accent_light"]
        # "restrained" — keep as-is (2 accent colors max)
        # "vibrant" — keep as-is, no restrictions

    def _luminance(self, hex_color):
        """Compute perceived luminance (0-255)."""
        h = hex_color.lstrip("#")
        if len(h) != 6:
            return 128
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return 0.299 * r + 0.587 * g + 0.114 * b

    def _is_dark(self, hex_color):
        h = hex_color.lstrip("#")
        if len(h) != 6:
            return False
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        luminance = 0.299 * r + 0.587 * g + 0.114 * b
        return luminance < 128

    def __getitem__(self, key):
        return self.d.get(key, DEFAULTS.get(key))

    def layout_for(self, layout_type):
        """Get learned layout pattern for a layout type, with fuzzy matching."""
        patterns = getattr(self, "_layout_patterns", {})
        if not patterns:
            return {}
        # Exact match
        if layout_type in patterns:
            return patterns[layout_type]
        # Fuzzy: match by lowercase keyword
        lt_lower = layout_type.lower()
        for name, pattern in patterns.items():
            if lt_lower in name.lower() or name.lower() in lt_lower:
                return pattern
        return {}


# ─── Slide Builders ────────────────────────────────────────────────────────

class SlideBuilder:
    """Builds individual slides on a Presentation object."""

    def __init__(self, prs, dl: DesignLanguage):
        self.prs = prs
        self.dl = dl
        self.slide_w = prs.slide_width
        self.slide_h = prs.slide_height
        self._img_temp_dir = tempfile.mkdtemp(prefix="executer_ppt_imgs_")

    def _pct_x(self, pct):
        return pct_to_emu(pct, self.slide_w)

    def _pct_y(self, pct):
        return pct_to_emu(pct, self.slide_h)

    def _add_slide(self):
        """Add a blank slide."""
        layout = self.prs.slide_layouts[6]  # blank layout
        return self.prs.slides.add_slide(layout)

    def _set_background(self, slide, color_hex=None):
        """Set slide background color."""
        bg_color = color_hex or self.dl["color_bg"]
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    def _add_textbox(self, slide, left, top, width, height,
                     text="", font_name=None, font_size=None,
                     bold=False, italic=False, color=None,
                     alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     line_spacing=None, auto_shrink=False):
        """Add a text box with styling. Auto-checks WCAG contrast."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_shrink else MSO_AUTO_SIZE.NONE

        try:
            tf.paragraphs[0].alignment = alignment
        except Exception:
            pass

        # Vertical alignment
        try:
            txBox.text_frame._txBody.bodyPr.set("anchor", {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }.get(anchor, "t"))
        except Exception:
            pass

        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = text

        actual_size = font_size or self.dl["body_size"]
        fg_color = color or self.dl["color_text"]
        bg_color = self.dl["color_bg"]

        # WCAG contrast check — large text (>=18pt bold or >=24pt) needs 3:1, else 4.5:1
        is_large_text = (actual_size >= 24) or (actual_size >= 18 and bold)
        min_ratio = 3.0 if is_large_text else 4.5
        fg_color = ensure_contrast(fg_color, bg_color, min_ratio)

        font = run.font
        font.name = font_name or self.dl["font_primary"]
        font.size = Pt(actual_size)
        font.bold = bold
        font.italic = italic
        font.color.rgb = hex_to_rgb(fg_color)

        if line_spacing:
            p.line_spacing = line_spacing

        return txBox

    def _add_bullets(self, slide, left, top, width, height,
                     bullets, font_name=None, font_size=None,
                     color=None, bold_first=False):
        """Add a bulleted list with WCAG contrast enforcement."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Pre-compute contrast-safe colors against slide background
        bg_color = self.dl["color_bg"]
        base_size = font_size or self.dl["bullet_size"]
        base_color = ensure_contrast(color or self.dl["color_body"], bg_color, 4.5)
        sub_color = ensure_contrast(color or self.dl["color_body"], bg_color, 4.5)

        # Apply max bullets from philosophy (content_density)
        max_bullets = getattr(self.dl, '_max_bullets', 6)
        visible_bullets = bullets[:max_bullets]

        first_para = True
        for i, bullet in enumerate(visible_bullets):
            # Handle nested bullets (sub-items) — level 1
            if isinstance(bullet, list):
                for j, sub in enumerate(bullet):
                    if not first_para:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                        first_para = False
                    # Handle triple-nested (level 2)
                    if isinstance(sub, list):
                        for k, subsub in enumerate(sub):
                            if k > 0:
                                p = tf.add_paragraph()
                            p.text = f"    \u25B8 {subsub}"
                            p.level = 2
                            run = p.runs[0]
                            run.font.name = font_name or self.dl["font_primary"]
                            run.font.size = Pt(base_size - 4)
                            run.font.color.rgb = hex_to_rgb(sub_color)
                            p.space_after = Pt(2)
                    else:
                        p.text = f"  \u2013 {sub}"
                        p.level = 1
                        run = p.runs[0]
                        run.font.name = font_name or self.dl["font_primary"]
                        run.font.size = Pt(base_size - 2)
                        run.font.color.rgb = hex_to_rgb(sub_color)
                        p.space_after = Pt(4)
                continue

            if not first_para:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
                first_para = False

            p.text = f"\u2022 {bullet}"
            p.level = 0
            run = p.runs[0]
            run.font.name = font_name or self.dl["font_primary"]
            run.font.size = Pt(base_size)
            run.font.color.rgb = hex_to_rgb(base_color)
            if bold_first and i == 0:
                run.font.bold = True
            p.space_after = Pt(6)
            p.line_spacing = self.dl["line_spacing"]

        return txBox

    # ── Visual Depth Primitives ──

    def _add_shadow(self, shape, blur_pt=None, offset_pt=None, color=None, alpha_pct=None):
        """Add a drop shadow to a shape via XML.
        Uses learned shadow parameters from the source deck if available,
        otherwise falls back to Apple-style subtle defaults."""
        try:
            from lxml import etree
            # Use learned shadow style from design language, with Apple-style fallbacks
            style = getattr(self.dl, '_shadow_style', {}) if hasattr(self.dl, '_shadow_style') else {}
            blur = blur_pt or style.get("blur_pt", 6)
            offset = offset_pt or style.get("offset_pt", 2)
            shd_color = color or style.get("color", "000000")
            alpha = alpha_pct or style.get("alpha_pct", 20)

            spPr = shape._element.spPr
            effectLst = spPr.find(qn('a:effectLst'))
            if effectLst is None:
                effectLst = etree.SubElement(spPr, qn('a:effectLst'))
            outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
            outerShdw.set('blurRad', str(int(blur * 12700)))
            outerShdw.set('dist', str(int(offset * 12700)))
            outerShdw.set('dir', '2700000')  # bottom-right
            outerShdw.set('algn', 'tl')
            outerShdw.set('rotWithShape', '0')
            srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
            srgbClr.set('val', shd_color.lstrip("#"))
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_elem.set('val', str(int(alpha * 1000)))
        except Exception as e:
            log(f"[warning] Shadow effect failed: {e}")

    def _set_transparency(self, shape, alpha_pct):
        """Set fill transparency on a shape (0=opaque, 100=invisible)."""
        try:
            from lxml import etree
            fill_elem = shape.fill._fill
            solidFill = fill_elem.find(qn('a:solidFill'))
            if solidFill is not None and len(solidFill) > 0:
                color_elem = solidFill[0]
                alpha_elem = etree.SubElement(color_elem, qn('a:alpha'))
                alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))
        except Exception as e:
            log(f"[warning] Transparency effect failed: {e}")

    def _add_rounded_rect(self, slide, left, top, width, height,
                          fill_color=None, shadow=None, border_color=None, border_width_pt=0):
        """Add a rounded rectangle with learned corner radius and shadow style."""
        # Use plain rectangle if source deck doesn't use rounded corners
        use_rounded = getattr(self.dl, '_use_rounded', False)
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if use_rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        else:
            shape.fill.background()
        if border_color:
            shape.line.color.rgb = hex_to_rgb(border_color)
            shape.line.width = Pt(border_width_pt or 1)
        else:
            shape.line.fill.background()
        # Apply learned corner radius from source deck
        if use_rounded and hasattr(self.dl, '_corner_radius_pct') and self.dl._corner_radius_pct:
            try:
                from lxml import etree
                prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
                if prstGeom is not None:
                    avLst = prstGeom.find(qn('a:avLst'))
                    if avLst is None:
                        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
                    # Clear existing guides and set learned radius
                    for child in list(avLst):
                        avLst.remove(child)
                    gd = etree.SubElement(avLst, qn('a:gd'))
                    gd.set('name', 'adj')
                    gd.set('fmla', f'val {int(self.dl._corner_radius_pct / 100 * 50000)}')
            except Exception as e:
                log(f"[warning] Rounded corner effect failed: {e}")
        # Shadow only when explicitly requested AND source deck uses shadows
        use_shadow = shadow if shadow is not None else getattr(self.dl, '_use_shadows', False)
        if use_shadow:
            self._add_shadow(shape)
        return shape

    def _add_circle(self, slide, left, top, size, fill_color, text=None,
                    text_color="FFFFFF", text_size=14):
        """Add a circle shape, optionally with centered text (WCAG contrast enforced)."""
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        shape.line.fill.background()
        if text:
            # Ensure text is readable against the circle fill color
            text_color = ensure_contrast(text_color, fill_color, 3.0)
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(text)
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(text_size)
                run.font.color.rgb = hex_to_rgb(text_color)
                run.font.bold = True
            try:
                tf._txBody.bodyPr.set("anchor", "ctr")
            except Exception:
                pass
        return shape

    def _add_card(self, slide, left, top, width, height,
                  title=None, body=None, bullets=None, icon_text=None,
                  fill_color=None, title_color=None):
        """Add a card component: rect with optional shadow + internal content."""
        card_fill = fill_color or self.dl["color_bg_subtle"]
        card_bg = self._add_rounded_rect(
            slide, left, top, width, height,
            fill_color=card_fill,
            border_color=self.dl["color_mid_gray"] if not getattr(self.dl, '_use_shadows', False) else None,
            border_width_pt=0.5 if not getattr(self.dl, '_use_shadows', False) else 0,
        )

        # Internal padding — scaled by whitespace philosophy
        pad_factor = getattr(self.dl, '_card_padding_factor', 1.0)
        pad_x = int(width * 0.08 * pad_factor)
        pad_y = int(height * 0.08 * pad_factor)
        inner_left = left + pad_x
        inner_w = width - 2 * pad_x
        content_top = top + pad_y

        # Optional icon circle at top
        if icon_text:
            circle_size = min(int(height * 0.16), int(Pt(32)))
            self._add_circle(
                slide, inner_left, content_top, circle_size,
                fill_color=self.dl["color_accent"],
                text=icon_text, text_size=12,
            )
            content_top += circle_size + int(Pt(6))

        # Title
        if title:
            th = int(height * 0.18)
            self._add_textbox(
                slide, inner_left, content_top, inner_w, th,
                text=title,
                font_size=self.dl["body_size"] + 2,
                bold=True,
                color=title_color or self.dl["color_text"],
            )
            content_top += th

        # Body or bullets
        remaining = max((top + height - pad_y) - content_top, int(Pt(10)))
        if body:
            self._add_textbox(
                slide, inner_left, content_top, inner_w, remaining,
                text=body,
                font_size=self.dl["body_size"] - 2,
                color=self.dl["color_body"],
                line_spacing=self.dl["line_spacing"],
                auto_shrink=True,
            )
        elif bullets:
            self._add_bullets(
                slide, inner_left + Pt(4), content_top, inner_w - Pt(8), remaining,
                bullets=bullets, font_size=self.dl["bullet_size"] - 2,
            )

        return card_bg

    def _add_slide_number(self, slide, number):
        """Add slide number to bottom-right."""
        self._add_textbox(
            slide,
            left=self._pct_x(90),
            top=self._pct_y(93),
            width=self._pct_x(8),
            height=self._pct_y(5),
            text=str(number),
            font_size=self.dl["note_size"],
            color=self.dl["color_slide_number"],
            alignment=PP_ALIGN.RIGHT,
        )

    def _add_accent_bar(self, slide, top_pct=18, width_pct=12, height_pt=4):
        """Add an accent bar. Respects learned accent_bar properties from design language."""
        # Check if accent bars were detected in the learned design
        ab = {}
        lp_all = getattr(self.dl, "_layout_patterns", {})
        for layout_name, pattern in (lp_all.items() if isinstance(lp_all, dict) else []):
            if "accent_bar" in pattern:
                ab = pattern["accent_bar"]
                break

        # If learned accent bar has very small dimensions, the source deck likely has no bars
        if ab and ab.get("width_pct", 10) < 1 and ab.get("height_pct", 1) < 0.3:
            return None  # Source deck has no accent bars

        left = self._pct_x(ab.get("left_pct", self.dl["margin_left_pct"]))
        top = self._pct_y(ab.get("top_pct", top_pct))
        width = self._pct_x(ab.get("width_pct", width_pct))
        height = Pt(height_pt)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self.dl["color_accent"])
        shape.line.fill.background()
        return shape

    # ── Layout: Title ──

    def build_title(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        lp = self.dl.layout_for("title")
        tp = lp.get("title", {})
        sp = lp.get("subtitle", {})

        title = content.get("title", "")
        subtitle = content.get("subtitle", "")

        # Title — centered, upper half (positions from learned layout or defaults)
        self._add_textbox(
            slide,
            left=self._pct_x(tp.get("left_pct", 10)),
            top=self._pct_y(tp.get("top_pct", 30)),
            width=self._pct_x(tp.get("width_pct", 80)),
            height=self._pct_y(tp.get("height_pct", 20)),
            text=title,
            font_size=self.dl["title_size"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.BOTTOM,
        )

        # Accent bar
        self._add_accent_bar(slide, top_pct=tp.get("top_pct", 30) + tp.get("height_pct", 20) + 3)

        # Subtitle
        if subtitle:
            self._add_textbox(
                slide,
                left=self._pct_x(sp.get("left_pct", 10)),
                top=self._pct_y(sp.get("top_pct", 57)),
                width=self._pct_x(sp.get("width_pct", 80)),
                height=self._pct_y(sp.get("height_pct", 12)),
                text=subtitle,
                font_size=self.dl["subtitle_size"],
                color=self.dl["color_text_secondary"],
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.TOP,
            )

        # No slide number on title slides
        return slide

    # ── Layout: Section ──

    def build_section(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        lp = self.dl.layout_for("section")
        tp = lp.get("title", {})
        sp = lp.get("subtitle", {})

        title = content.get("title", "")
        subtitle = content.get("subtitle", "")

        # Centered large title
        self._add_textbox(
            slide,
            left=self._pct_x(tp.get("left_pct", 10)),
            top=self._pct_y(tp.get("top_pct", 35)),
            width=self._pct_x(tp.get("width_pct", 80)),
            height=self._pct_y(tp.get("height_pct", 20)),
            text=title,
            font_size=self.dl["title_size"] * 0.9,
            bold=True,
            color=self.dl["color_accent"],
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.BOTTOM,
        )

        if subtitle:
            self._add_textbox(
                slide,
                left=self._pct_x(sp.get("left_pct", 15)),
                top=self._pct_y(sp.get("top_pct", 58)),
                width=self._pct_x(sp.get("width_pct", 70)),
                height=self._pct_y(sp.get("height_pct", 10)),
                text=subtitle,
                font_size=self.dl["subtitle_size"],
                color=self.dl["color_text_secondary"],
                alignment=PP_ALIGN.CENTER,
            )

        # No slide number on section slides
        return slide

    # ── Layout: Content ──

    def build_content(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        # Title
        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title,
                font_size=self.dl["heading_size"],
                bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        # Check for image
        image = content.get("image", None)
        img_position = "right"  # default
        if isinstance(image, dict):
            img_position = image.get("position", "right")
        elif isinstance(image, str) and image:
            image = {"url": image}

        # Body text
        body = content.get("body", "")
        bullets = content.get("bullets", [])
        content_top = self.dl["content_top_pct"] + 5

        if image and img_position == "right":
            # Split layout: text left, image right
            text_w = content_w * 0.55
            img_w_pct = content_w * 0.40
            img_left = margin_l + text_w + content_w * 0.05
            img_h_pct = image.get("height_pct", 55) if isinstance(image, dict) else 55
            img_w_pct_final = image.get("width_pct", img_w_pct) if isinstance(image, dict) else img_w_pct

            self._resolve_and_place_image(
                slide, image,
                self._pct_x(img_left), self._pct_y(content_top),
                self._pct_x(img_w_pct_final), self._pct_y(img_h_pct),
            )
        elif image and img_position == "left":
            text_w = content_w * 0.55
            img_w_pct = content_w * 0.40
            img_h_pct = image.get("height_pct", 55) if isinstance(image, dict) else 55
            img_w_pct_final = image.get("width_pct", img_w_pct) if isinstance(image, dict) else img_w_pct

            self._resolve_and_place_image(
                slide, image,
                self._pct_x(margin_l), self._pct_y(content_top),
                self._pct_x(img_w_pct_final), self._pct_y(img_h_pct),
            )
        else:
            text_w = content_w

        # Determine text left position
        if image and img_position == "left":
            txt_left = margin_l + (image.get("width_pct", content_w * 0.40) if isinstance(image, dict) else content_w * 0.40) + content_w * 0.05
        else:
            txt_left = margin_l

        actual_text_w = text_w if image and img_position in ("left", "right") else content_w

        if body:
            self._add_textbox(
                slide,
                left=self._pct_x(txt_left), top=self._pct_y(content_top),
                width=self._pct_x(actual_text_w), height=self._pct_y(15),
                text=body,
                font_size=self.dl["body_size"],
                color=self.dl["color_body"],
                line_spacing=self.dl["line_spacing"],
            )
            content_top += 18

        if bullets:
            self._add_bullets(
                slide,
                left=self._pct_x(txt_left + 2), top=self._pct_y(content_top),
                width=self._pct_x(actual_text_w - 4), height=self._pct_y(55),
                bullets=bullets,
            )

        # Image at bottom (after text/bullets)
        if image and img_position == "bottom":
            img_h_pct = image.get("height_pct", 35) if isinstance(image, dict) else 35
            img_w_pct = image.get("width_pct", content_w * 0.6) if isinstance(image, dict) else content_w * 0.6
            img_top = 88 - img_h_pct - 2
            img_left = margin_l + (content_w - img_w_pct) / 2  # center
            self._resolve_and_place_image(
                slide, image,
                self._pct_x(img_left), self._pct_y(img_top),
                self._pct_x(img_w_pct), self._pct_y(img_h_pct),
            )

        # Full background image
        if image and img_position == "full":
            self._resolve_and_place_image(
                slide, image,
                self._pct_x(margin_l), self._pct_y(content_top),
                self._pct_x(content_w), self._pct_y(85 - content_top),
            )

        # Note/footnote
        note = content.get("note", "")
        if note:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(88),
                width=self._pct_x(content_w), height=self._pct_y(5),
                text=note,
                font_size=self.dl["note_size"],
                color=self.dl["color_slide_number"],
                italic=True,
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Two Column ──

    def build_two_column(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        # Title
        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title,
                font_size=self.dl["heading_size"],
                bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        col_top = self.dl["content_top_pct"] + 5
        left_start = margin_l
        col_w = (content_w - 6) / 2  # 6% gap
        right_start = margin_l + col_w + 6

        # Divider line
        divider_x = self._pct_x(50)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            divider_x, self._pct_y(col_top),
            Pt(1), self._pct_y(60)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.dl["color_mid_gray"])
        line.line.fill.background()

        # Left column
        left_title = content.get("left_title", "")
        left_bullets = content.get("left_bullets", [])
        left_image = content.get("left_image", None)
        if left_title:
            self._add_textbox(
                slide,
                left=self._pct_x(left_start), top=self._pct_y(col_top),
                width=self._pct_x(col_w), height=self._pct_y(7),
                text=left_title,
                font_size=self.dl["body_size"] + 2,
                bold=True,
                color=self.dl["color_accent"],
            )
        if left_image:
            img_top = col_top + (9 if left_title else 0)
            img_h = 50 if not left_bullets else 25
            self._resolve_and_place_image(
                slide, left_image,
                self._pct_x(left_start), self._pct_y(img_top),
                self._pct_x(col_w), self._pct_y(img_h),
            )
            if left_bullets:
                bullet_top = img_top + img_h + 2
            else:
                bullet_top = None
        else:
            bullet_top = col_top + 9 if left_bullets else None
        if left_bullets and bullet_top is not None:
            self._add_bullets(
                slide,
                left=self._pct_x(left_start + 1), top=self._pct_y(bullet_top),
                width=self._pct_x(col_w - 2), height=self._pct_y(50 if not left_image else 22),
                bullets=left_bullets,
            )

        # Right column
        right_title = content.get("right_title", "")
        right_bullets = content.get("right_bullets", [])
        right_image = content.get("right_image", None)
        if right_title:
            self._add_textbox(
                slide,
                left=self._pct_x(right_start), top=self._pct_y(col_top),
                width=self._pct_x(col_w), height=self._pct_y(7),
                text=right_title,
                font_size=self.dl["body_size"] + 2,
                bold=True,
                color=self.dl["color_accent2"],
            )
        if right_image:
            img_top = col_top + (9 if right_title else 0)
            img_h = 50 if not right_bullets else 25
            self._resolve_and_place_image(
                slide, right_image,
                self._pct_x(right_start), self._pct_y(img_top),
                self._pct_x(col_w), self._pct_y(img_h),
            )
            if right_bullets:
                bullet_top_r = img_top + img_h + 2
            else:
                bullet_top_r = None
        else:
            bullet_top_r = col_top + 9 if right_bullets else None
        if right_bullets and bullet_top_r is not None:
            self._add_bullets(
                slide,
                left=self._pct_x(right_start + 1), top=self._pct_y(bullet_top_r),
                width=self._pct_x(col_w - 2), height=self._pct_y(50 if not right_image else 22),
                bullets=right_bullets,
            )

        note = content.get("note", "")
        if note:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(88),
                width=self._pct_x(content_w), height=self._pct_y(5),
                text=note,
                font_size=self.dl["note_size"],
                color=self.dl["color_slide_number"],
                italic=True,
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Quote ──

    def build_quote(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        quote = content.get("quote_text", content.get("body", ""))
        attribution = content.get("attribution", "")

        # Large opening quote mark
        self._add_textbox(
            slide,
            left=self._pct_x(8), top=self._pct_y(20),
            width=self._pct_x(10), height=self._pct_y(15),
            text="\u201C",
            font_size=72,
            color=self.dl["color_accent"],
            bold=True,
        )

        # Quote text
        quote_size = min(self.dl["body_size"] * 1.4, 32)
        self._add_textbox(
            slide,
            left=self._pct_x(12), top=self._pct_y(30),
            width=self._pct_x(76), height=self._pct_y(35),
            text=quote,
            font_size=quote_size,
            italic=True,
            color=self.dl["color_text"],
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.4,
        )

        # Closing quote mark
        self._add_textbox(
            slide,
            left=self._pct_x(82), top=self._pct_y(55),
            width=self._pct_x(10), height=self._pct_y(15),
            text="\u201D",
            font_size=72,
            color=self.dl["color_accent"],
            bold=True,
        )

        # Attribution
        if attribution:
            self._add_textbox(
                slide,
                left=self._pct_x(20), top=self._pct_y(72),
                width=self._pct_x(60), height=self._pct_y(7),
                text=attribution,
                font_size=self.dl["body_size"],
                color=self.dl["color_text_secondary"],
                alignment=PP_ALIGN.CENTER,
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Image Left / Image Right ──

    def _build_image_slide(self, content, slide_num, image_side="left"):
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]
        img_source = content.get("image_url") or content.get("image_path", "")
        title = content.get("title", "")
        body = content.get("body", "")
        bullets = content.get("bullets", [])

        # Title full width
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title,
                font_size=self.dl["heading_size"],
                bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        col_top = self.dl["content_top_pct"] + 5
        col_w = (content_w - 4) / 2

        if image_side == "left":
            img_left = margin_l
            text_left = margin_l + col_w + 4
        else:
            text_left = margin_l
            img_left = margin_l + col_w + 4

        # Image or placeholder
        img_x = self._pct_x(img_left)
        img_y = self._pct_y(col_top)
        img_w = self._pct_x(col_w)
        img_h = self._pct_y(60)

        self._resolve_and_place_image(slide, img_source, img_x, img_y, img_w, img_h)

        # Text side
        if body:
            self._add_textbox(
                slide,
                left=self._pct_x(text_left), top=self._pct_y(col_top),
                width=self._pct_x(col_w), height=self._pct_y(15),
                text=body,
                font_size=self.dl["body_size"],
                color=self.dl["color_body"],
                line_spacing=self.dl["line_spacing"],
            )

        if bullets:
            bullet_top = col_top + (18 if body else 0)
            self._add_bullets(
                slide,
                left=self._pct_x(text_left + 1), top=self._pct_y(bullet_top),
                width=self._pct_x(col_w - 2), height=self._pct_y(45),
                bullets=bullets,
            )

        self._add_slide_number(slide, slide_num)
        return slide

    def build_image_left(self, content, slide_num):
        return self._build_image_slide(content, slide_num, "left")

    def build_image_right(self, content, slide_num):
        return self._build_image_slide(content, slide_num, "right")

    def _add_placeholder_rect(self, slide, x, y, w, h, label=""):
        """Add a placeholder rectangle where an image would go."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self.dl["color_light_gray"])
        shape.line.color.rgb = hex_to_rgb(self.dl["color_mid_gray"])
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        display = Path(label).name if label else "Image"
        p.text = f"[{display}]"
        run = p.runs[0]
        run.font.size = Pt(self.dl["caption_size"])
        run.font.color.rgb = hex_to_rgb(self.dl["color_slide_number"])

    def _resolve_and_place_image(self, slide, img_spec, x, y, w, h):
        """Resolve an image (path or URL) and place it on the slide.

        img_spec can be:
          - A string (path or URL)
          - A dict with 'url', 'path', and optional size overrides
        Returns True if image was placed, False if skipped.
        On failure: fills the area with a subtle accent block instead of an
        ugly gray placeholder — keeps the slide looking clean.
        """
        if isinstance(img_spec, dict):
            source = img_spec.get("url") or img_spec.get("path", "")
        else:
            source = str(img_spec) if img_spec else ""

        if not source:
            self._add_clean_image_fallback(slide, x, y, w, h)
            return False

        local_path = resolve_image(source, self._img_temp_dir)
        if local_path and Path(local_path).exists():
            try:
                slide.shapes.add_picture(local_path, x, y, w, h)
                return True
            except Exception:
                self._add_clean_image_fallback(slide, x, y, w, h)
                return False
        else:
            self._add_clean_image_fallback(slide, x, y, w, h)
            return False

    def _add_clean_image_fallback(self, slide, x, y, w, h):
        """When an image fails to load, add a subtle tinted rectangle
        instead of a gray placeholder with a broken filename.
        Keeps the slide looking intentionally designed."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self.dl["color_bg_subtle"])
        shape.line.fill.background()  # No border — clean look

    # ── Layout: Data Table ──

    def build_data_table(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title,
                font_size=self.dl["heading_size"],
                bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        headers = content.get("headers", [])
        rows = content.get("rows", [])
        if not headers and not rows:
            self._add_slide_number(slide, slide_num)
            return slide

        num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
        num_rows = (1 if headers else 0) + len(rows)

        table_top = self.dl["content_top_pct"] + 5
        table_x = self._pct_x(margin_l)
        table_y = self._pct_y(table_top)
        table_w = self._pct_x(content_w)
        table_h = self._pct_y(min(num_rows * 8, 60))

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, table_x, table_y, table_w, table_h
        )
        table = table_shape.table

        # Style header row
        row_idx = 0
        if headers:
            for col_idx, header_text in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header_text)

                # Header background — use accent color, auto-pick text color for contrast
                header_bg = self.dl["color_accent"]
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = hex_to_rgb(header_bg)

                # WCAG contrast: ensure header text is readable on accent background
                header_text_hex = ensure_contrast("FFFFFF", header_bg, 4.5)
                # If white doesn't pass (very light accent), try dark text
                if wcag_contrast_ratio("FFFFFF", header_bg) < 4.5:
                    header_text_hex = ensure_contrast(self.dl["color_text"], header_bg, 4.5)
                header_text_color = hex_to_rgb(header_text_hex)

                # Header text style
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(self.dl["body_size"] - 2)
                        run.font.color.rgb = header_text_color
                        run.font.name = self.dl["font_primary"]

                # Cell margins
                self._set_cell_margins(cell, Inches(0.1))
            row_idx = 1

        # Data rows with alternating colors
        for data_row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    break
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)

                # Alternating row color (uses accent-derived tint for color coherence)
                if data_row_idx % 2 == 1:
                    cell_fill = cell.fill
                    cell_fill.solid()
                    cell_fill.fore_color.rgb = hex_to_rgb(self.dl["color_bg_subtle"])

                # WCAG contrast for data cells against their row background
                cell_bg = self.dl["color_bg_subtle"] if data_row_idx % 2 == 1 else self.dl["color_bg"]
                cell_text_color = ensure_contrast(self.dl["color_body"], cell_bg, 4.5)
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(self.dl["body_size"] - 2)
                        run.font.name = self.dl["font_primary"]
                        run.font.color.rgb = hex_to_rgb(cell_text_color)

                self._set_cell_margins(cell, Inches(0.1))
            row_idx += 1

        # Smart column widths — proportional to max content length
        col_max_lens = [len(str(h)) for h in headers] if headers else [5] * num_cols
        for row_data in rows:
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < num_cols:
                    col_max_lens[col_idx] = max(col_max_lens[col_idx], len(str(cell_text)))
        total_len = sum(col_max_lens) or 1
        for col_idx, col in enumerate(table.columns):
            col.width = int(table_w * col_max_lens[col_idx] / total_len)

        note = content.get("note", "")
        if note:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(88),
                width=self._pct_x(content_w), height=self._pct_y(5),
                text=note,
                font_size=self.dl["note_size"],
                color=self.dl["color_slide_number"],
                italic=True,
            )

        self._add_slide_number(slide, slide_num)
        return slide

    def _set_cell_margins(self, cell, margin):
        """Set cell margins via XML (python-pptx doesn't expose this directly)."""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcPr.set("marL", str(int(margin)))
        tcPr.set("marR", str(int(margin)))
        tcPr.set("marT", str(int(margin / 2)))
        tcPr.set("marB", str(int(margin / 2)))

    # ── Layout: Blank ──

    def build_blank(self, content, slide_num):
        slide = self._add_slide()
        self._set_background(slide)

        # Support arbitrary image placement on blank slides
        images = content.get("images", [])
        for img_spec in images:
            if isinstance(img_spec, str):
                img_spec = {"url": img_spec, "left_pct": 10, "top_pct": 10, "width_pct": 80, "height_pct": 80}
            left_pct = img_spec.get("left_pct", 10)
            top_pct = img_spec.get("top_pct", 10)
            width_pct = img_spec.get("width_pct", 80)
            height_pct = img_spec.get("height_pct", 80)
            self._resolve_and_place_image(
                slide, img_spec,
                self._pct_x(left_pct), self._pct_y(top_pct),
                self._pct_x(width_pct), self._pct_y(height_pct),
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Big Number ──

    def build_big_number(self, content, slide_num):
        """Large statistic/metric with supporting text."""
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        number = content.get("number", content.get("title", ""))
        label = content.get("label", content.get("subtitle", ""))
        body = content.get("body", "")

        # Large number — centered, prominent
        number_size = min(self.dl["title_size"] * 2.5, 120)
        self._add_textbox(
            slide,
            left=self._pct_x(margin_l), top=self._pct_y(25),
            width=self._pct_x(content_w), height=self._pct_y(30),
            text=number,
            font_size=number_size,
            bold=True,
            color=self.dl["color_accent"],
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.BOTTOM,
        )

        # Label
        if label:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(58),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=label,
                font_size=self.dl["heading_size"],
                color=self.dl["color_text_secondary"],
                alignment=PP_ALIGN.CENTER,
            )

        # Supporting body text
        if body:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l + 10), top=self._pct_y(72),
                width=self._pct_x(content_w - 20), height=self._pct_y(15),
                text=body,
                font_size=self.dl["body_size"],
                color=self.dl["color_body"],
                alignment=PP_ALIGN.CENTER,
                line_spacing=self.dl["line_spacing"],
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Full Image ──

    def build_full_image(self, content, slide_num):
        """Full-bleed background image with text overlay."""
        slide = self._add_slide()
        self._set_background(slide)

        image = content.get("image", content.get("image_url", ""))
        title = content.get("title", "")
        subtitle = content.get("subtitle", content.get("body", ""))

        # Full-bleed image
        if image:
            self._resolve_and_place_image(
                slide, image,
                0, 0,
                self.slide_w, self.slide_h,
            )

        # Semi-transparent overlay bar for text readability
        if title or subtitle:
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, self._pct_y(60),
                self.slide_w, self._pct_y(35),
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            self._set_transparency(overlay, 50)  # 50% see-through
            overlay.line.fill.background()

            if title:
                self._add_textbox(
                    slide,
                    left=self._pct_x(8), top=self._pct_y(65),
                    width=self._pct_x(84), height=self._pct_y(12),
                    text=title,
                    font_size=self.dl["heading_size"],
                    bold=True,
                    color="FFFFFF",
                    alignment=PP_ALIGN.LEFT,
                )
            if subtitle:
                self._add_textbox(
                    slide,
                    left=self._pct_x(8), top=self._pct_y(78),
                    width=self._pct_x(84), height=self._pct_y(10),
                    text=subtitle,
                    font_size=self.dl["body_size"],
                    color="EEEEEE",
                    alignment=PP_ALIGN.LEFT,
                )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Comparison ──

    def build_comparison(self, content, slide_num):
        """Side-by-side comparison with visual emphasis."""
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title,
                font_size=self.dl["heading_size"],
                bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        col_top = self.dl["content_top_pct"] + 5
        col_w = (content_w - 8) / 2  # wider gap than two_column
        left_start = margin_l
        right_start = margin_l + col_w + 8

        # Left panel background
        self._add_rounded_rect(
            slide,
            self._pct_x(left_start - 1), self._pct_y(col_top - 1),
            self._pct_x(col_w + 2), self._pct_y(65),
            fill_color=self.dl["color_bg_subtle"],
        )

        # Left content
        left_title = content.get("left_title", "")
        left_bullets = content.get("left_bullets", [])
        if left_title:
            self._add_textbox(
                slide,
                left=self._pct_x(left_start + 1), top=self._pct_y(col_top + 2),
                width=self._pct_x(col_w - 2), height=self._pct_y(8),
                text=left_title,
                font_size=self.dl["body_size"] + 4,
                bold=True,
                color=self.dl["color_accent"],
                alignment=PP_ALIGN.CENTER,
            )
        if left_bullets:
            self._add_bullets(
                slide,
                left=self._pct_x(left_start + 3), top=self._pct_y(col_top + 12),
                width=self._pct_x(col_w - 6), height=self._pct_y(48),
                bullets=left_bullets,
            )

        # "VS" divider
        self._add_textbox(
            slide,
            left=self._pct_x(margin_l + col_w + 1), top=self._pct_y(col_top + 25),
            width=self._pct_x(6), height=self._pct_y(8),
            text="VS",
            font_size=self.dl["body_size"] + 2,
            bold=True,
            color=self.dl["color_text_secondary"],
            alignment=PP_ALIGN.CENTER,
        )

        # Right content
        right_title = content.get("right_title", "")
        right_bullets = content.get("right_bullets", [])
        if right_title:
            self._add_textbox(
                slide,
                left=self._pct_x(right_start + 1), top=self._pct_y(col_top + 2),
                width=self._pct_x(col_w - 2), height=self._pct_y(8),
                text=right_title,
                font_size=self.dl["body_size"] + 4,
                bold=True,
                color=self.dl["color_accent2"],
                alignment=PP_ALIGN.CENTER,
            )
        if right_bullets:
            self._add_bullets(
                slide,
                left=self._pct_x(right_start + 3), top=self._pct_y(col_top + 12),
                width=self._pct_x(col_w - 6), height=self._pct_y(48),
                bullets=right_bullets,
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Cards ──

    def build_cards(self, content, slide_num):
        """Grid of styled cards — modern feature/value display."""
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title, font_size=self.dl["heading_size"], bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        cards = content.get("cards", [])
        if not cards:
            self._add_slide_number(slide, slide_num)
            return slide

        n = len(cards)
        card_top = self.dl["content_top_pct"] + 5
        card_area_h = max(68 - card_top, 15)  # leave room for slide number

        # layout_complexity limits grid density
        complexity = getattr(self.dl, '_layout_complexity', 'moderate')
        max_cols = 2 if complexity == "simple" else (4 if complexity == "complex" else 3)

        if n <= 2:
            cols, rows = min(n, max_cols), 1
        elif n <= 3:
            cols, rows = min(n, max_cols), 1
        elif n <= 4:
            cols, rows = min(2, max_cols), 2
        elif n <= 6:
            cols, rows = min(3, max_cols), 2
        else:
            cols, rows = min(4, max_cols), 2

        # Whitespace style affects card gap
        ws = getattr(self.dl, '_whitespace_style', 'balanced')
        gap_pct = 5 if ws == "generous" else (2 if ws == "tight" else 3)
        card_w_pct = (content_w - gap_pct * (cols - 1)) / cols
        card_h_pct = (card_area_h - gap_pct * (rows - 1)) / rows

        for idx, card_data in enumerate(cards[:cols * rows]):
            row = idx // cols
            col = idx % cols
            cx = margin_l + col * (card_w_pct + gap_pct)
            cy = card_top + row * (card_h_pct + gap_pct)

            self._add_card(
                slide,
                self._pct_x(cx), self._pct_y(cy),
                self._pct_x(card_w_pct), self._pct_y(card_h_pct),
                title=card_data.get("title"),
                body=card_data.get("body"),
                bullets=card_data.get("bullets"),
                icon_text=card_data.get("icon"),
            )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Process Flow ──

    def build_process(self, content, slide_num):
        """Horizontal process flow with numbered circles and connectors."""
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title, font_size=self.dl["heading_size"], bold=True,
            )
            self._add_accent_bar(slide, top_pct=self.dl["margin_top_pct"] + 10)

        steps = content.get("steps", [])
        if not steps:
            self._add_slide_number(slide, slide_num)
            return slide

        n = len(steps)
        step_top = self.dl["content_top_pct"] + 8
        circle_size_pct = min(10, content_w / n * 0.4)
        circle_size = self._pct_x(circle_size_pct)

        # Evenly space steps
        total_gap = content_w - circle_size_pct * n
        spacing = total_gap / max(n - 1, 1) if n > 1 else 0

        for i, step in enumerate(steps):
            cx_pct = margin_l + i * (circle_size_pct + spacing)
            cx = self._pct_x(cx_pct)
            cy = self._pct_y(step_top)

            # Connector line (between circles, not before first)
            if i > 0:
                prev_cx = self._pct_x(margin_l + (i - 1) * (circle_size_pct + spacing))
                line_left = prev_cx + circle_size
                line_width = cx - line_left
                if line_width > 0:
                    connector = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        line_left, cy + int(circle_size * 0.45),
                        line_width, Pt(3),
                    )
                    connector.fill.solid()
                    connector.fill.fore_color.rgb = hex_to_rgb(self.dl["color_accent_light"])
                    connector.line.fill.background()

            # Circle with step number
            self._add_circle(
                slide, cx, cy, circle_size,
                fill_color=self.dl["color_accent"],
                text=str(i + 1), text_size=max(10, int(circle_size_pct * 1.5)),
            )

            # Scale label/description widths to available space
            slot_w = circle_size_pct + (spacing if n > 1 else content_w * 0.3)
            label_w = min(slot_w, circle_size_pct + 4)
            desc_w = min(slot_w, circle_size_pct + 8)
            label_offset = (label_w - circle_size_pct) / 2
            desc_offset = (desc_w - circle_size_pct) / 2

            # Label below circle
            label = step.get("label", "")
            if label:
                self._add_textbox(
                    slide,
                    left=self._pct_x(cx_pct - label_offset), top=self._pct_y(step_top + circle_size_pct + 2),
                    width=self._pct_x(label_w), height=self._pct_y(7),
                    text=label,
                    font_size=self.dl["body_size"],
                    bold=True,
                    alignment=PP_ALIGN.CENTER,
                )

            # Description below label
            desc = step.get("description", "")
            if desc:
                self._add_textbox(
                    slide,
                    left=self._pct_x(cx_pct - desc_offset), top=self._pct_y(step_top + circle_size_pct + 10),
                    width=self._pct_x(desc_w), height=self._pct_y(12),
                    text=desc,
                    font_size=self.dl["body_size"] - 3,
                    color=self.dl["color_text_secondary"],
                    alignment=PP_ALIGN.CENTER,
                    line_spacing=self.dl["line_spacing"],
                )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Layout: Timeline ──

    def build_timeline(self, content, slide_num):
        """Horizontal timeline with date markers and alternating event descriptions."""
        slide = self._add_slide()
        self._set_background(slide)

        margin_l = self.dl["margin_left_pct"]
        content_w = 100 - margin_l - self.dl["margin_right_pct"]

        title = content.get("title", "")
        if title:
            self._add_textbox(
                slide,
                left=self._pct_x(margin_l), top=self._pct_y(self.dl["margin_top_pct"]),
                width=self._pct_x(content_w), height=self._pct_y(10),
                text=title, font_size=self.dl["heading_size"], bold=True,
            )

        events = content.get("events", [])
        if not events:
            self._add_slide_number(slide, slide_num)
            return slide

        n = len(events)
        line_y_pct = 55  # horizontal line at vertical center-ish

        # Main timeline line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._pct_x(margin_l), self._pct_y(line_y_pct),
            self._pct_x(content_w), Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.dl["color_accent"])
        line.line.fill.background()

        # Events along the line — scale label widths to available space
        event_spacing = content_w / max(n, 1)
        dot_size = self._pct_x(2)
        label_w = min(14, event_spacing * 1.4)
        desc_w = min(16, event_spacing * 1.6)

        for i, event in enumerate(events):
            ex_pct = margin_l + event_spacing * (i + 0.5) - 1
            ex = self._pct_x(ex_pct)

            # Dot on timeline
            self._add_circle(
                slide, ex, self._pct_y(line_y_pct) - int(dot_size * 0.3),
                dot_size,
                fill_color=self.dl["color_accent"],
            )

            # Alternate above/below
            above = (i % 2 == 0)

            date = event.get("date", "")
            label = event.get("label", "")
            desc = event.get("description", "")

            if above:
                # Date above line
                if date:
                    self._add_textbox(
                        slide,
                        left=self._pct_x(ex_pct - label_w / 2), top=self._pct_y(line_y_pct - 22),
                        width=self._pct_x(label_w), height=self._pct_y(6),
                        text=date, font_size=self.dl["caption_size"],
                        color=self.dl["color_accent"], bold=True,
                        alignment=PP_ALIGN.CENTER,
                    )
                if label:
                    self._add_textbox(
                        slide,
                        left=self._pct_x(ex_pct - label_w / 2), top=self._pct_y(line_y_pct - 16),
                        width=self._pct_x(label_w), height=self._pct_y(7),
                        text=label, font_size=self.dl["body_size"] - 2,
                        bold=True, alignment=PP_ALIGN.CENTER,
                    )
                if desc:
                    self._add_textbox(
                        slide,
                        left=self._pct_x(ex_pct - desc_w / 2), top=self._pct_y(line_y_pct - 9),
                        width=self._pct_x(desc_w), height=self._pct_y(8),
                        text=desc, font_size=self.dl["caption_size"],
                        color=self.dl["color_text_secondary"],
                        alignment=PP_ALIGN.CENTER,
                    )
            else:
                # Below line
                if date:
                    self._add_textbox(
                        slide,
                        left=self._pct_x(ex_pct - label_w / 2), top=self._pct_y(line_y_pct + 4),
                        width=self._pct_x(label_w), height=self._pct_y(6),
                        text=date, font_size=self.dl["caption_size"],
                        color=self.dl["color_accent"], bold=True,
                        alignment=PP_ALIGN.CENTER,
                    )
                if label:
                    self._add_textbox(
                        slide,
                        left=self._pct_x(ex_pct - label_w / 2), top=self._pct_y(line_y_pct + 10),
                        width=self._pct_x(label_w), height=self._pct_y(7),
                        text=label, font_size=self.dl["body_size"] - 2,
                        bold=True, alignment=PP_ALIGN.CENTER,
                    )
                if desc:
                    self._add_textbox(
                        slide,
                        left=self._pct_x(ex_pct - desc_w / 2), top=self._pct_y(line_y_pct + 17),
                        width=self._pct_x(desc_w), height=self._pct_y(8),
                        text=desc, font_size=self.dl["caption_size"],
                        color=self.dl["color_text_secondary"],
                        alignment=PP_ALIGN.CENTER,
                    )

        self._add_slide_number(slide, slide_num)
        return slide

    # ── Dispatcher ──

    def build_slide(self, slide_spec, slide_num):
        """Route to the correct layout builder."""
        layout = slide_spec.get("layout", "content")
        content = slide_spec.get("content", {})

        builders = {
            "title": self.build_title,
            "section": self.build_section,
            "content": self.build_content,
            "two_column": self.build_two_column,
            "quote": self.build_quote,
            "image_left": self.build_image_left,
            "image_right": self.build_image_right,
            "data_table": self.build_data_table,
            "blank": self.build_blank,
            "big_number": self.build_big_number,
            "full_image": self.build_full_image,
            "comparison": self.build_comparison,
            "cards": self.build_cards,
            "process": self.build_process,
            "timeline": self.build_timeline,
        }

        builder = builders.get(layout, self.build_content)
        slide = builder(content, slide_num)

        # Speaker notes support (applies to all layouts)
        notes = content.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide


# ─── Main ──────────────────────────────────────────────────────────────────

# ─── Spec Advisor ─────────────────────────────────────────────────────────

def advise_and_fix_spec(spec):
    """Analyze a slide spec and auto-fix common problems.

    Runs BEFORE the engine renders. Returns (fixed_spec, warnings).
    This is the intelligence layer between the LLM and the dumb renderer.
    """
    warnings = []
    slides = spec.get("slides", [])
    if not slides:
        return spec, ["No slides in spec"]

    fixed_slides = []
    for i, slide in enumerate(slides):
        layout = slide.get("layout", "content")
        content = slide.get("content", {})

        # ── Rule 1: Split dense content slides ──
        # If a slide with bullets has 7+, split into multiple slides
        if layout in ("content", "image_left", "image_right") and len(content.get("bullets", [])) > 6:
            bullets = content["bullets"]
            title = content.get("title", "")
            # Split into chunks of 4
            for chunk_idx in range(0, len(bullets), 4):
                chunk = bullets[chunk_idx:chunk_idx + 4]
                chunk_title = title if chunk_idx == 0 else f"{title} (cont.)"
                new_content = dict(content)
                new_content["bullets"] = chunk
                new_content["title"] = chunk_title
                fixed_slides.append({"layout": layout, "content": new_content})
            warnings.append(f"Slide {i+1}: Split {len(bullets)} bullets into {(len(bullets)+3)//4} slides")
            continue

        # ── Rule 2: Truncate excessively long titles ──
        title = content.get("title", "")
        if len(title) > 80:
            content["title"] = title[:77] + "..."
            warnings.append(f"Slide {i+1}: Truncated title from {len(title)} chars")

        # ── Rule 3: Truncate long bullet text ──
        bullets = content.get("bullets", [])
        fixed_bullets = []
        for b in bullets:
            if isinstance(b, str) and len(b) > 120:
                fixed_bullets.append(b[:117] + "...")
                warnings.append(f"Slide {i+1}: Truncated a bullet from {len(b)} chars")
            else:
                fixed_bullets.append(b)
        if bullets:
            content["bullets"] = fixed_bullets

        # ── Rule 4: Auto-generate speaker notes if missing ──
        if not content.get("notes"):
            notes = _generate_speaker_notes(layout, content)
            if notes:
                content["notes"] = notes

        fixed_slides.append({"layout": layout, "content": content})

    # ── Rule 5: Enforce layout rhythm ──
    # Flag 4+ identical layouts in a row
    if len(fixed_slides) >= 4:
        run_count = 1
        for j in range(1, len(fixed_slides)):
            if fixed_slides[j]["layout"] == fixed_slides[j-1]["layout"]:
                run_count += 1
                if run_count >= 4:
                    warnings.append(
                        f"Slides {j-2}-{j+1}: {run_count} identical '{fixed_slides[j]['layout']}' "
                        f"layouts in a row — consider adding a section divider or varying the layout"
                    )
            else:
                run_count = 1

    # ── Rule 6: Warn on layout monotony ──
    content_layouts = [s["layout"] for s in fixed_slides if s["layout"] not in ("title", "section")]
    if len(content_layouts) >= 3:
        unique = set(content_layouts)
        if len(unique) == 1:
            warnings.append(
                f"All {len(content_layouts)} content slides use '{content_layouts[0]}' layout — "
                f"vary with cards, big_number, comparison, process, image_right, or full_image"
            )

    # ── Rule 7: Strip hallucinated/invalid image references ──
    # Only keep image URLs that look like real URLs (http/https) or existing local paths.
    # LLMs frequently hallucinate filenames like "ai-trends-2025.jpg" that don't exist.
    import os as _os
    image_keys = ("image", "image_url", "image_path", "left_image", "right_image",
                  "background_image")
    stripped_count = 0
    for s in fixed_slides:
        c = s.get("content", {})
        for key in image_keys:
            val = c.get(key)
            if not val:
                continue
            src = val.get("url", val) if isinstance(val, dict) else val
            if isinstance(src, str) and src:
                is_url = src.startswith("http://") or src.startswith("https://")
                is_local = _os.path.isabs(src) and _os.path.exists(src)
                if not is_url and not is_local:
                    # Hallucinated filename — strip it
                    c.pop(key, None)
                    stripped_count += 1
        # Downgrade image-only layouts to content if image was stripped
        if s["layout"] in ("full_image", "image_left", "image_right"):
            has_any_img = any(c.get(k) for k in image_keys)
            if not has_any_img:
                s["layout"] = "content"
                stripped_count += 1
    if stripped_count > 0:
        warnings.append(
            f"Stripped {stripped_count} hallucinated image references (not real URLs or files). "
            f"Use search_images FIRST to get real image URLs before creating the spec."
        )

    # Warn on imageless decks (after stripping)
    if len(fixed_slides) >= 5:
        has_image = False
        for s in fixed_slides:
            c = s.get("content", {})
            if s["layout"] in ("full_image", "image_left", "image_right"):
                has_image = True
                break
            if any(c.get(k) for k in image_keys):
                has_image = True
                break
        if not has_image:
            warnings.append(
                f"No images in a {len(fixed_slides)}-slide deck — use search_images to find "
                f"relevant visuals, then add full_image or image_right slides"
            )

    # ── Rule 8: Suggest cards for short bullet lists ──
    for i, s in enumerate(fixed_slides):
        if s["layout"] in ("content", "image_left", "image_right"):
            bullets = s.get("content", {}).get("bullets", [])
            if 3 <= len(bullets) <= 5:
                avg_len = sum(len(str(b)) for b in bullets if isinstance(b, str)) / max(len(bullets), 1)
                if avg_len < 30:
                    warnings.append(
                        f"Slide {i+1}: {len(bullets)} short bullets (avg {int(avg_len)} chars) — "
                        f"consider 'cards' layout for cleaner visual presentation"
                    )

    spec["slides"] = fixed_slides
    return spec, warnings


def _generate_speaker_notes(layout, content):
    """Generate contextual speaker notes based on slide content."""
    title = content.get("title", "")
    bullets = content.get("bullets", [])
    body = content.get("body", "")

    if layout == "title":
        subtitle = content.get("subtitle", "")
        if subtitle:
            return f"Welcome everyone. Today's topic: {title}. {subtitle}."
        return f"Welcome everyone. Today we'll discuss: {title}."

    if layout == "section":
        return f"Let's move on to the next section: {title}."

    if layout == "big_number":
        number = content.get("number", "")
        label = content.get("label", "")
        return f"The key number here is {number} — {label}. {body}" if body else f"The key takeaway: {number} — {label}."

    if layout == "content" and bullets:
        flat = [b for b in bullets if isinstance(b, str)]
        if flat:
            return f"On this slide about {title}: " + "; ".join(flat[:3]) + ("..." if len(flat) > 3 else ".")

    if layout == "comparison":
        left = content.get("left_title", "")
        right = content.get("right_title", "")
        if left and right:
            return f"Here we compare {left} versus {right}."

    if layout == "quote":
        quote = content.get("quote_text", "")
        attr = content.get("attribution", "")
        if quote:
            return f"As {attr} said: \"{quote[:100]}{'...' if len(quote) > 100 else ''}\""

    if layout in ("process", "timeline"):
        steps = content.get("steps", content.get("events", []))
        labels = [s.get("label", "") for s in steps if s.get("label")]
        if labels:
            return f"{title}: " + " → ".join(labels) + "."

    if layout == "cards":
        cards = content.get("cards", [])
        titles = [c.get("title", "") for c in cards if c.get("title")]
        if titles:
            return f"Let's look at each: " + ", ".join(titles) + "."

    return None


def create_ppt(spec, design_path=None, output_path="output.pptx"):
    """Create a .pptx from spec dict + optional design language."""
    try:
        # Run spec advisor — fix problems BEFORE rendering
        spec, advisor_warnings = advise_and_fix_spec(spec)
        for w in advisor_warnings:
            log(f"[advisor] {w}")

        dl = DesignLanguage(design_path)

        prs = Presentation()
        prs.slide_width = Inches(dl["slide_width"]) if isinstance(dl["slide_width"], (int, float)) else dl["slide_width"]
        prs.slide_height = Inches(dl["slide_height"]) if isinstance(dl["slide_height"], (int, float)) else dl["slide_height"]

        builder = SlideBuilder(prs, dl)

        # Log design philosophy if available
        if hasattr(dl, '_content_density'):
            log(f"Design: density={dl._content_density}, whitespace={dl._whitespace_style}, "
                f"complexity={dl._layout_complexity}, shadows={'yes' if dl._use_shadows else 'no'}, "
                f"gradients={'yes' if dl._use_gradients else 'no'}")

        slides = spec.get("slides", [])
        if not slides:
            return {"success": False, "error": "No slides in spec"}

        for i, slide_spec in enumerate(slides):
            slide_num = i + 1
            log(f"Building slide {slide_num}/{len(slides)}: {slide_spec.get('layout', 'content')}")
            builder.build_slide(slide_spec, slide_num)

        # Ensure output directory exists
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)

        prs.save(str(out))
        log(f"Saved: {out}")

        # Clean up temp downloaded images
        if builder._img_temp_dir:
            cleanup_temp_images(builder._img_temp_dir)

        result = {
            "success": True,
            "path": str(out.resolve()),
            "slides": len(slides),
            "design_language": design_path is not None,
        }
        if advisor_warnings:
            result["advisor_notes"] = advisor_warnings
        return result

    except Exception as e:
        return {"success": False, "error": f"{type(e).__name__}: {str(e)}"}


def main():
    parser = argparse.ArgumentParser(description="Executer PPT Engine")
    parser.add_argument("--spec", help="Path to spec JSON file")
    parser.add_argument("--stdin", action="store_true", help="Read spec from stdin")
    parser.add_argument("--design", help="Path to design_language.json")
    parser.add_argument("--output", required=True, help="Output .pptx path")

    args = parser.parse_args()

    # Load spec
    if args.stdin:
        try:
            spec = json.load(sys.stdin)
        except json.JSONDecodeError as e:
            print(json.dumps({"success": False, "error": f"Invalid JSON from stdin: {e}"}))
            sys.exit(1)
    elif args.spec:
        try:
            spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
        except FileNotFoundError:
            print(json.dumps({"success": False, "error": f"Spec file not found: {args.spec}"}))
            sys.exit(1)
        except json.JSONDecodeError as e:
            print(json.dumps({"success": False, "error": f"Invalid spec JSON: {e}"}))
            sys.exit(1)
    else:
        print(json.dumps({"success": False, "error": "Provide --spec or --stdin"}))
        sys.exit(1)

    # Handle filename from spec
    output = args.output
    if output.endswith("/") or Path(output).is_dir():
        filename = spec.get("filename", "presentation.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        output = str(Path(output) / filename)

    result = create_ppt(spec, design_path=args.design, output_path=output)
    print(json.dumps(result))
    sys.exit(0 if result["success"] else 1)


if __name__ == "__main__":
    main()
