#!/usr/bin/env python3
"""
PPT Design Language Extractor
Extracts the complete design system from a .pptx file:
  - Slide dimensions & layout structure
  - Color palette (every color used, with frequency)
  - Typography system (fonts, sizes, weights, line spacing)
  - Element positioning grid (where things are placed)
  - Shape catalog (types, sizes, fill patterns)
  - Text hierarchy (title vs subtitle vs body patterns)
  - Image placement patterns
  - Slide-by-slide structural blueprint

Usage:
    python ppt_design_extractor.py input.pptx
    python ppt_design_extractor.py input.pptx -o design_language.json
    python ppt_design_extractor.py input.pptx --format markdown
"""

import argparse
import json
import sys
from collections import Counter, defaultdict
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("Install python-pptx first:")
    print("  pip install python-pptx")
    sys.exit(1)


def emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    if emu is None:
        return None
    return round(emu / 914400, 2)


def emu_to_pt(emu):
    """Convert EMU to points."""
    if emu is None:
        return None
    return round(emu / 12700, 1)


def extract_color(color_obj):
    """Safely extract color as hex string."""
    try:
        if color_obj is None:
            return None
        if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
            return f"#{color_obj.rgb}"
        if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
            return f"theme:{color_obj.theme_color}"
    except Exception:
        pass
    return None


def extract_fill(shape):
    """Extract fill information from a shape."""
    try:
        fill = shape.fill
        if fill is None:
            return None
        fill_type = str(fill.type) if fill.type is not None else None
        result = {"type": fill_type}
        if fill.type is not None:
            ft = str(fill.type)
            if "SOLID" in ft:
                try:
                    result["color"] = extract_color(fill.fore_color)
                except Exception:
                    pass
            elif "GRADIENT" in ft:
                stops = []
                try:
                    for stop in fill.gradient_stops:
                        stops.append({
                            "position": round(stop.position, 2),
                            "color": extract_color(stop.color)
                        })
                except Exception:
                    pass
                result["gradient_stops"] = stops
                # Extract gradient angle/direction from XML
                try:
                    from lxml import etree
                    gradFill = shape._element.spPr.find(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
                    if gradFill is not None:
                        lin = gradFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
                        if lin is not None:
                            angle = lin.get('ang')
                            if angle:
                                result["gradient_angle_deg"] = int(angle) // 60000
                except Exception:
                    pass
        return result
    except Exception:
        return None


def extract_text_formatting(paragraph):
    """Extract detailed text formatting from a paragraph."""
    runs_info = []
    for run in paragraph.runs:
        run_data = {"text": run.text}
        font = run.font
        if font.name:
            run_data["font"] = font.name
        if font.size:
            run_data["size_pt"] = emu_to_pt(font.size)
        if font.bold:
            run_data["bold"] = True
        if font.italic:
            run_data["italic"] = True
        if font.underline:
            run_data["underline"] = True
        color = extract_color(font.color)
        if color:
            run_data["color"] = color
        runs_info.append(run_data)

    para_data = {}
    if paragraph.alignment is not None:
        align_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
        }
        para_data["alignment"] = align_map.get(paragraph.alignment, str(paragraph.alignment))

    # paragraph_format may not exist in all python-pptx versions — access safely
    try:
        pf = paragraph.paragraph_format
        if pf.space_before is not None:
            para_data["space_before_pt"] = emu_to_pt(pf.space_before)
        if pf.space_after is not None:
            para_data["space_after_pt"] = emu_to_pt(pf.space_after)
        if pf.line_spacing is not None:
            para_data["line_spacing"] = round(pf.line_spacing, 2) if isinstance(pf.line_spacing, float) else emu_to_pt(pf.line_spacing)
        if pf.level is not None and pf.level > 0:
            para_data["indent_level"] = pf.level
    except AttributeError:
        # Fallback: try direct XML access for spacing
        try:
            from lxml import etree
            pPr = paragraph._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is not None:
                lvl = pPr.get('lvl')
                if lvl and int(lvl) > 0:
                    para_data["indent_level"] = int(lvl)
                spcBef = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef')
                if spcBef is not None:
                    spcPts = spcBef.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
                    if spcPts is not None:
                        para_data["space_before_pt"] = int(spcPts.get('val', '0')) / 100
                spcAft = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft')
                if spcAft is not None:
                    spcPts = spcAft.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
                    if spcPts is not None:
                        para_data["space_after_pt"] = int(spcPts.get('val', '0')) / 100
        except Exception:
            pass

    if runs_info:
        para_data["runs"] = runs_info

    return para_data


def extract_shape_info(shape, slide_width, slide_height):
    """Extract complete information from a shape."""
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else None,
    }

    # Position & size in inches + as percentage of slide
    if shape.left is not None:
        info["position"] = {
            "left_in": emu_to_inches(shape.left),
            "top_in": emu_to_inches(shape.top),
            "width_in": emu_to_inches(shape.width),
            "height_in": emu_to_inches(shape.height),
            "left_pct": round(shape.left / slide_width * 100, 1) if slide_width else None,
            "top_pct": round(shape.top / slide_height * 100, 1) if slide_height else None,
            "width_pct": round(shape.width / slide_width * 100, 1) if slide_width else None,
            "height_pct": round(shape.height / slide_height * 100, 1) if slide_height else None,
        }

    # Rotation
    if hasattr(shape, 'rotation') and shape.rotation:
        info["rotation_deg"] = shape.rotation

    # Fill
    fill = extract_fill(shape)
    if fill:
        info["fill"] = fill

    # Line/border
    try:
        line = shape.line
        if line and line.fill and line.fill.type is not None:
            line_info = {}
            if line.width:
                line_info["width_pt"] = emu_to_pt(line.width)
            line_color = extract_color(line.color)
            if line_color:
                line_info["color"] = line_color
            if line.dash_style:
                line_info["dash"] = str(line.dash_style)
            if line_info:
                info["border"] = line_info
    except Exception:
        pass

    # Shadow — extract full parameters, not just boolean
    try:
        if hasattr(shape, 'shadow') and shape.shadow:
            shadow = shape.shadow
            if shadow.inherit is False:
                shadow_info = {"has_shadow": True}
                # Extract shadow parameters from XML for full fidelity
                try:
                    from lxml import etree
                    spPr = shape._element.spPr
                    effectLst = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
                    if effectLst is not None:
                        outerShdw = effectLst.find('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw')
                        if outerShdw is not None:
                            blur = outerShdw.get('blurRad')
                            if blur:
                                shadow_info["blur_pt"] = round(int(blur) / 12700, 1)
                            dist = outerShdw.get('dist')
                            if dist:
                                shadow_info["offset_pt"] = round(int(dist) / 12700, 1)
                            direction = outerShdw.get('dir')
                            if direction:
                                shadow_info["direction_deg"] = int(direction) // 60000
                            # Shadow color + alpha
                            srgb = outerShdw.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                            if srgb is not None:
                                shadow_info["color"] = srgb.get('val', '000000')
                                alpha = srgb.find('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                                if alpha is not None:
                                    shadow_info["alpha_pct"] = round(int(alpha.get('val', '100000')) / 1000, 1)
                except Exception:
                    pass
                info["shadow"] = shadow_info
    except Exception:
        pass

    # Corner radius (adjustment values for rounded rectangles)
    try:
        from lxml import etree
        prstGeom = shape._element.spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is not None and prstGeom.get('prst') == 'roundRect':
            avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
            if avLst is not None:
                for gd in avLst.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                    if gd.get('name') == 'adj':
                        # Value is in 1/50000 of shape size
                        try:
                            raw = int(gd.get('fmla', 'val 16667').split()[-1])
                            info["corner_radius_pct"] = round(raw / 50000 * 100, 1)
                        except (ValueError, IndexError):
                            info["corner_radius_pct"] = 33.3  # default roundRect radius
    except Exception:
        pass

    # Element opacity/transparency
    try:
        from lxml import etree
        fill_elem = shape._element.spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if fill_elem is not None and len(fill_elem) > 0:
            alpha = fill_elem[0].find('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            if alpha is not None:
                info["opacity_pct"] = round(int(alpha.get('val', '100000')) / 1000, 1)
    except Exception:
        pass

    # Text content
    if shape.has_text_frame:
        tf = shape.text_frame
        info["text_frame"] = {
            "word_wrap": tf.word_wrap,
        }
        # Margins
        margins = {}
        if tf.margin_left is not None:
            margins["left_in"] = emu_to_inches(tf.margin_left)
        if tf.margin_right is not None:
            margins["right_in"] = emu_to_inches(tf.margin_right)
        if tf.margin_top is not None:
            margins["top_in"] = emu_to_inches(tf.margin_top)
        if tf.margin_bottom is not None:
            margins["bottom_in"] = emu_to_inches(tf.margin_bottom)
        if margins:
            info["text_frame"]["margins"] = margins

        # Paragraphs
        paragraphs = []
        for para in tf.paragraphs:
            pdata = extract_text_formatting(para)
            if pdata:
                paragraphs.append(pdata)
        if paragraphs:
            info["text_frame"]["paragraphs"] = paragraphs

    # Table
    if shape.has_table:
        table = shape.table
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "col_widths_in": [emu_to_inches(col.width) for col in table.columns],
            "row_heights_in": [emu_to_inches(row.height) for row in table.rows],
        }
        # Extract first row as header sample
        if len(table.rows) > 0:
            header_cells = []
            for cell in table.rows[0].cells:
                cell_data = {"text": cell.text}
                if cell.fill and cell.fill.type is not None:
                    cell_data["fill"] = extract_fill(cell)
                header_cells.append(cell_data)
            info["table"]["header_sample"] = header_cells

    # Chart
    if shape.has_chart:
        chart = shape.chart
        info["chart"] = {
            "chart_type": str(chart.chart_type),
            "has_legend": chart.has_legend,
            "series_count": len(chart.series) if chart.series else 0,
        }

    # Image / Picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            info["image"] = {
                "content_type": img.content_type,
                "width_px": img.size[0] if hasattr(img, 'size') else None,
                "height_px": img.size[1] if hasattr(img, 'size') else None,
            }
        except Exception:
            info["image"] = {"note": "embedded image"}

    # Group
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["group_children"] = len(shape.shapes)

    # Placeholder info
    if shape.is_placeholder:
        ph = shape.placeholder_format
        info["placeholder"] = {
            "idx": ph.idx,
            "type": str(ph.type),
        }

    return info


def extract_slide_layout_info(layout):
    """Extract info about a slide layout template."""
    info = {
        "name": layout.name,
        "placeholders": []
    }
    for ph in layout.placeholders:
        info["placeholders"].append({
            "idx": ph.placeholder_format.idx,
            "type": str(ph.placeholder_format.type),
            "name": ph.name,
            "position": {
                "left_in": emu_to_inches(ph.left),
                "top_in": emu_to_inches(ph.top),
                "width_in": emu_to_inches(ph.width),
                "height_in": emu_to_inches(ph.height),
            }
        })
    return info


def _hex_luminance(hex_str):
    """Compute perceived luminance (0-255) from a hex color string."""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        return 128
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _is_near_white(hex_str, threshold=230):
    return _hex_luminance(hex_str) >= threshold


def _is_near_black(hex_str, threshold=40):
    return _hex_luminance(hex_str) <= threshold


def analyze_design_system(slides_data, all_colors, all_fonts, all_sizes,
                          text_colors=None, fill_colors=None,
                          bg_colors=None, border_colors=None,
                          theme_scheme=None):
    """Derive the design system from collected data with semantic color classification."""
    system = {}

    # Color palette — sorted by frequency (kept for backwards compatibility)
    color_counts = Counter(c for c in all_colors if c and not c.startswith("theme:"))
    if color_counts:
        system["color_palette"] = [
            {"hex": color, "uses": count}
            for color, count in color_counts.most_common(20)
        ]

    # Theme colors
    theme_colors = Counter(c for c in all_colors if c and c.startswith("theme:"))
    if theme_colors:
        system["theme_colors"] = [
            {"ref": color, "uses": count}
            for color, count in theme_colors.most_common()
        ]

    # === Semantic color classification ===
    # Strategy: detect whether deck is light-mode or dark-mode first,
    # then assign text/bg/accent roles accordingly.
    semantic = {}

    # Collect all non-theme colors
    tc = Counter(c for c in (text_colors or []) if c and not c.startswith("theme:"))
    bc = Counter(c for c in (bg_colors or []) if c and not c.startswith("theme:"))
    fc_all = Counter(c for c in (fill_colors or []) if c and not c.startswith("theme:"))
    all_notheme = Counter(c for c in all_colors if c and not c.startswith("theme:"))

    # Detect dark mode: use overall color distribution weighted by frequency.
    # If the most-used colors are dark, it's a dark deck — even if bg is inherited.
    top_all = all_notheme.most_common(10)
    if top_all:
        # Weight by frequency — the dominant visual impression
        total_uses = sum(n for _, n in top_all)
        avg_lum = sum(_hex_luminance(c) * n for c, n in top_all) / max(total_uses, 1)
    else:
        avg_lum = 200  # assume light if no colors

    is_dark_mode = avg_lum < 100

    if is_dark_mode:
        # Dark deck: darkest frequent fill = background, lightest text = text
        dark_fills = sorted(fc_all.most_common(10), key=lambda x: _hex_luminance(x[0]))
        if dark_fills:
            semantic["background"] = dark_fills[0][0]
        # Text: pick the lightest frequent text color
        if tc:
            light_texts = sorted(tc.most_common(10), key=lambda x: -_hex_luminance(x[0]))
            semantic["text_primary"] = light_texts[0][0]
            if len(light_texts) >= 2:
                semantic["text_secondary"] = light_texts[1][0]
    else:
        # Light deck: lightest fill = background, darkest text = text
        if bc:
            semantic["background"] = bc.most_common(1)[0][0]
        elif fc_all:
            light_fills = [(c, n) for c, n in fc_all.most_common(10) if _hex_luminance(c) > 180]
            if light_fills:
                semantic["background"] = light_fills[0][0]
        if tc:
            semantic["text_primary"] = tc.most_common(1)[0][0]
            if len(tc) >= 2:
                semantic["text_secondary"] = tc.most_common(2)[1][0]

    # Accent colors: the most vivid colors that aren't background or near-neutral
    bg_hex = semantic.get("background", "").lstrip("#").upper()
    text_hex = semantic.get("text_primary", "").lstrip("#").upper()

    # Score by saturation — vivid colors make the best accents
    def _color_saturation(hex_str):
        h = hex_str.lstrip("#")
        if len(h) != 6: return 0
        r, g, b = int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16)
        mx, mn = max(r, g, b), min(r, g, b)
        return (mx - mn) / max(mx, 1) * 255

    accent_candidates = [
        (c, n, _color_saturation(c)) for c, n in all_notheme.most_common(30)
        if c.lstrip("#").upper() != bg_hex
        and c.lstrip("#").upper() != text_hex
        and not _is_near_white(c)
        and not _is_near_black(c)
    ]
    # Sort by saturation * frequency — vivid and common wins
    accent_candidates.sort(key=lambda x: x[2] * x[1], reverse=True)

    if accent_candidates:
        semantic["accent"] = accent_candidates[0][0]
        if len(accent_candidates) >= 2:
            semantic["accent2"] = accent_candidates[1][0]

    # If we still don't have an accent (monochrome deck), derive from text
    if "accent" not in semantic and text_hex:
        semantic["accent"] = semantic.get("text_primary", "#333333")

    # === Theme-based fallback ===
    # When the PPT uses mostly theme colors, the RGB-based analysis above can be
    # unreliable (e.g., white fills counted as text). The theme's dk1/lt1/accent1
    # directly encode the designer's intent and are authoritative.
    if theme_scheme:
        # Always use theme for text/bg if they look wrong (text shouldn't be white on white bg)
        if semantic.get("text_primary") and semantic.get("background"):
            tp = semantic["text_primary"].lstrip("#").upper()
            bg = semantic["background"].lstrip("#").upper()
            if tp == bg:
                # Text and background are the same — use theme
                if "dk1" in theme_scheme:
                    semantic["text_primary"] = theme_scheme["dk1"]
                if "lt1" in theme_scheme:
                    semantic["background"] = theme_scheme["lt1"]
        # Fill in missing values from theme
        if "text_primary" not in semantic and "dk1" in theme_scheme:
            semantic["text_primary"] = theme_scheme["dk1"]
        if "background" not in semantic and "lt1" in theme_scheme:
            semantic["background"] = theme_scheme["lt1"]
        if ("accent" not in semantic or _is_near_white(semantic.get("accent", "")) or
                _is_near_black(semantic.get("accent", ""))) and "accent1" in theme_scheme:
            semantic["accent"] = theme_scheme["accent1"]
        if "accent2" not in semantic and "accent2" in theme_scheme:
            semantic["accent2"] = theme_scheme["accent2"]
        if "text_secondary" not in semantic and "dk2" in theme_scheme:
            semantic["text_secondary"] = theme_scheme["dk2"]

    if semantic:
        system["semantic_colors"] = semantic

    # Typography
    font_counts = Counter(f for f in all_fonts if f)
    size_counts = Counter(s for s in all_sizes if s)
    system["typography"] = {
        "fonts_by_frequency": [
            {"font": font, "uses": count}
            for font, count in font_counts.most_common()
        ],
        "sizes_by_frequency": [
            {"size_pt": size, "uses": count}
            for size, count in size_counts.most_common()
        ],
    }

    # Detect text hierarchy
    if size_counts:
        sorted_sizes = sorted(size_counts.keys(), reverse=True)
        hierarchy = []
        role_names = ["title", "subtitle", "heading", "subheading", "body", "caption", "footnote"]
        for i, size in enumerate(sorted_sizes[:len(role_names)]):
            hierarchy.append({
                "likely_role": role_names[i],
                "size_pt": size,
                "uses": size_counts[size]
            })
        system["text_hierarchy"] = hierarchy

    return system


def derive_layout_patterns(slides_data):
    """Derive spatial layout patterns from per-slide shape data.

    Groups slides by layout name and computes median positions for
    identified shape roles (title, body, image, accent bar, etc.).
    Also derives global margin and spacing metrics.
    """
    from statistics import median

    layout_groups = {}
    for slide in slides_data:
        layout = slide.get("layout_name") or "Unknown"
        layout_groups.setdefault(layout, []).append(slide)

    layout_patterns = {}
    all_left_pcts = []
    all_right_edges = []
    all_top_pcts = []

    for layout_name, slides in layout_groups.items():
        roles = {}  # role -> list of position dicts

        for slide in slides:
            shapes = slide.get("shapes", [])
            # Sort shapes by top position to identify roles
            positioned = [(s, s.get("position", {})) for s in shapes if s.get("position")]
            if not positioned:
                continue

            # Identify title: top-most text shape with large font
            text_shapes = []
            for s, pos in positioned:
                tf = s.get("text_frame", {})
                paras = tf.get("paragraphs", [])
                if paras:
                    max_size = 0
                    for p in paras:
                        for r in p.get("runs", []):
                            sz = r.get("size_pt", 0)
                            if sz > max_size:
                                max_size = sz
                    if max_size > 0:
                        text_shapes.append((s, pos, max_size))

            # Sort by font size desc — largest is likely title
            text_shapes.sort(key=lambda x: x[2], reverse=True)

            if text_shapes:
                # Title: largest text
                _, tpos, _ = text_shapes[0]
                roles.setdefault("title", []).append(tpos)
                all_left_pcts.append(tpos.get("left_pct") or 5)
                right_edge = (tpos.get("left_pct") or 0) + (tpos.get("width_pct") or 0)
                all_right_edges.append(right_edge)
                all_top_pcts.append(tpos.get("top_pct") or 12)

                # Subtitle: second largest text (if notably smaller than title)
                if len(text_shapes) >= 2 and text_shapes[1][2] < text_shapes[0][2] * 0.85:
                    _, spos, _ = text_shapes[1]
                    roles.setdefault("subtitle", []).append(spos)

                # Body area: remaining text shapes
                for _, bpos, sz in text_shapes[2:]:
                    roles.setdefault("body_area", []).append(bpos)
                    all_left_pcts.append(bpos.get("left_pct") or 5)

            # Detect accent bars: thin shapes with fill, no text
            for s, pos in positioned:
                h = pos.get("height_pct", 100)
                w = pos.get("width_pct", 100)
                has_text = bool(s.get("text_frame", {}).get("paragraphs", []))
                has_fill = bool(s.get("fill"))

                if has_fill and not has_text:
                    if h is not None and h < 2 and w is not None and w > 5:  # Horizontal bar
                        roles.setdefault("accent_bar", []).append(pos)
                    elif w is not None and w < 2 and h is not None and h > 5:  # Vertical bar
                        roles.setdefault("accent_bar", []).append(pos)

            # Detect images
            for s, pos in positioned:
                if s.get("image"):
                    roles.setdefault("image", []).append(pos)

        # Compute median position for each role
        if roles:
            pattern = {}
            for role, positions in roles.items():
                if positions:
                    pattern[role] = {
                        "left_pct": round(median(p.get("left_pct", 0) or 0 for p in positions), 1),
                        "top_pct": round(median(p.get("top_pct", 0) or 0 for p in positions), 1),
                        "width_pct": round(median(p.get("width_pct", 0) or 0 for p in positions), 1),
                        "height_pct": round(median(p.get("height_pct", 0) or 0 for p in positions), 1),
                    }
            if pattern:
                layout_patterns[layout_name] = pattern

    # Derive global spacing metrics
    global_spacing = {}
    if all_left_pcts:
        global_spacing["margin_left_pct"] = round(median(all_left_pcts), 1)
    if all_right_edges:
        global_spacing["margin_right_pct"] = round(100 - median(all_right_edges), 1)
    if all_top_pcts:
        global_spacing["margin_top_pct"] = round(median(all_top_pcts), 1)
        # Content typically starts after title + some gap
        global_spacing["content_top_pct"] = round(median(all_top_pcts) + 10, 1)

    return {"layout_patterns": layout_patterns, "global_spacing": global_spacing}


def extract_ppt(filepath):
    """Main extraction function."""
    prs = Presentation(filepath)

    # Slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    result = {
        "source_file": str(Path(filepath).name),
        "slide_dimensions": {
            "width_in": emu_to_inches(slide_width),
            "height_in": emu_to_inches(slide_height),
            "aspect_ratio": f"{round(slide_width/slide_height, 2)}:1" if slide_height else None,
            "width_emu": slide_width,
            "height_emu": slide_height,
        },
        "total_slides": len(prs.slides),
    }

    # Slide layouts available
    layouts = []
    for layout in prs.slide_layouts:
        layouts.append(extract_slide_layout_info(layout))
    result["available_layouts"] = layouts

    # Slide master colors (theme) — extract from theme part, not master element
    theme_info = {}
    try:
        master = prs.slide_masters[0]
        theme_info["name"] = master.name if hasattr(master, 'name') else None

        # Method 1: look for clrScheme in master element tree
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        theme_el = master.element.find(f'.//{ns_a}clrScheme')

        # Method 2: if not found, look in the theme part (linked via relationship)
        if theme_el is None:
            try:
                for rel in master.part.rels.values():
                    if 'theme' in str(rel.reltype).lower():
                        theme_blob = rel.target_part.blob
                        from lxml import etree
                        theme_root = etree.fromstring(theme_blob)
                        theme_el = theme_root.find(f'.//{ns_a}clrScheme')
                        if theme_el is not None:
                            theme_info["name"] = theme_el.get("name", theme_info.get("name"))
                        break
            except Exception:
                pass

        if theme_el is not None:
            theme_colors = {}
            for child in theme_el:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                for sub in child:
                    if 'srgbClr' in sub.tag:
                        theme_colors[tag] = f"#{sub.attrib.get('val', '000000')}"
                    elif 'sysClr' in sub.tag:
                        # System colors: use lastClr (actual rendered color)
                        theme_colors[tag] = f"#{sub.attrib.get('lastClr', sub.attrib.get('val', '000000'))}"
            if theme_colors:
                theme_info["color_scheme"] = theme_colors
                print(f"Theme colors resolved: {list(theme_colors.keys())}", file=sys.stderr)
        result["theme"] = theme_info
    except Exception as e:
        print(f"Warning: Theme extraction failed: {e}", file=sys.stderr)
        result["theme"] = theme_info

    # Build theme color resolution map: theme:NAME → #RRGGBB
    # Maps MSO_THEME_COLOR enum names to actual hex values from the theme
    _theme_resolve = {}
    tc_scheme = theme_info.get("color_scheme", {}) if "theme" in result else {}
    # Map python-pptx theme_color enum names to theme XML tag names
    _theme_tag_map = {
        "DARK_1": "dk1", "LIGHT_1": "lt1", "DARK_2": "dk2", "LIGHT_2": "lt2",
        "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
        "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
        "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
        "TEXT_1": "dk1", "TEXT_2": "dk2",           # TEXT maps to DK in Office themes
        "BACKGROUND_1": "lt1", "BACKGROUND_2": "lt2",  # BG maps to LT
    }
    for enum_name, tag in _theme_tag_map.items():
        if tag in tc_scheme:
            _theme_resolve[enum_name] = tc_scheme[tag]

    def resolve_color(color_str):
        """Resolve theme: prefixed colors to actual hex values."""
        if not color_str or not color_str.startswith("theme:"):
            return color_str
        # Extract the enum name: "theme:TEXT_1 (13)" → "TEXT_1"
        theme_ref = color_str.replace("theme:", "").strip()
        # Remove parenthetical enum value
        if " (" in theme_ref:
            theme_ref = theme_ref.split(" (")[0]
        resolved = _theme_resolve.get(theme_ref)
        if resolved:
            return resolved
        return None  # Can't resolve — drop it

    # Process each slide — track colors by semantic category
    all_colors = []
    text_colors = []
    fill_colors = []
    bg_colors = []
    border_colors = []
    all_fonts = []
    all_sizes = []
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": slide_idx + 1,
            "layout_name": slide.slide_layout.name if slide.slide_layout else None,
        }

        # Background
        try:
            bg = slide.background
            if bg.fill and bg.fill.type is not None:
                slide_info["background"] = extract_fill(bg)
                bg_color = extract_color(bg.fill.fore_color) if "SOLID" in str(bg.fill.type) else None
                if bg_color:
                    all_colors.append(bg_color)
                    bg_colors.append(bg_color)
        except Exception:
            pass

        # Shapes
        shapes = []
        for shape in slide.shapes:
            shape_info = extract_shape_info(shape, slide_width, slide_height)
            shapes.append(shape_info)

            # Collect colors by category for semantic analysis
            fill = shape_info.get("fill", {})
            if isinstance(fill, dict) and fill.get("color"):
                c = fill["color"]
                all_colors.append(c)
                fill_colors.append(c)

            border = shape_info.get("border", {})
            if isinstance(border, dict) and border.get("color"):
                c = border["color"]
                all_colors.append(c)
                border_colors.append(c)

            tf = shape_info.get("text_frame", {})
            for para in tf.get("paragraphs", []):
                for run in para.get("runs", []):
                    if run.get("font"):
                        all_fonts.append(run["font"])
                    if run.get("size_pt"):
                        all_sizes.append(run["size_pt"])
                    if run.get("color"):
                        c = run["color"]
                        all_colors.append(c)
                        text_colors.append(c)

        slide_info["shapes"] = shapes
        slide_info["shape_count"] = len(shapes)
        slides_data.append(slide_info)

    result["slides"] = slides_data

    # Resolve theme colors to RGB before analysis
    # Theme-only PPTs (e.g., Apple-style black/white with accent) would otherwise
    # produce 0 colors and fall back to ugly defaults.
    all_colors = [resolve_color(c) or c for c in all_colors]
    text_colors = [resolve_color(c) or c for c in text_colors]
    fill_colors = [resolve_color(c) or c for c in fill_colors]
    bg_colors = [resolve_color(c) or c for c in bg_colors]
    border_colors = [resolve_color(c) or c for c in border_colors]
    # Strip any remaining unresolved theme: refs
    all_colors = [c for c in all_colors if c and not c.startswith("theme:")]
    text_colors = [c for c in text_colors if c and not c.startswith("theme:")]
    fill_colors = [c for c in fill_colors if c and not c.startswith("theme:")]
    bg_colors = [c for c in bg_colors if c and not c.startswith("theme:")]
    border_colors = [c for c in border_colors if c and not c.startswith("theme:")]

    # Derive design system with semantic color classification
    result["design_system"] = analyze_design_system(
        slides_data, all_colors, all_fonts, all_sizes,
        text_colors=text_colors, fill_colors=fill_colors,
        bg_colors=bg_colors, border_colors=border_colors,
        theme_scheme=tc_scheme
    )

    # Structural patterns
    layout_usage = Counter(s["layout_name"] for s in slides_data if s.get("layout_name"))
    result["layout_usage"] = [
        {"layout": name, "slides_using": count}
        for name, count in layout_usage.most_common()
    ]

    # Shape type frequency
    shape_types = Counter()
    for s in slides_data:
        for sh in s.get("shapes", []):
            st = sh.get("shape_type", "unknown")
            shape_types[st] += 1
    result["shape_type_frequency"] = [
        {"type": t, "count": c} for t, c in shape_types.most_common()
    ]

    # Derive spatial layout patterns and global spacing
    spatial = derive_layout_patterns(slides_data)
    result["layout_patterns"] = spatial.get("layout_patterns", {})
    result["global_spacing"] = spatial.get("global_spacing", {})

    # Derive visual effects summary — now captures PARAMETERS, not just booleans
    visual_effects = {
        "has_shadows": False,
        "shadow_count": 0,
        "shadow_style": None,          # Aggregated: typical shadow params from source
        "has_gradients": False,
        "gradient_count": 0,
        "gradient_angles": [],         # Collected angles for typical direction
        "has_rounded_corners": False,
        "corner_radius_pct": None,     # Typical corner radius
        "has_transparency": False,
        "shape_variety": [],
        "decorative_elements": [],
    }
    shadow_params = []
    corner_radii = []
    gradient_angles = []
    for slide in slides_data:
        for shape in slide.get("shapes", []):
            # Shadow with full parameters
            shadow = shape.get("shadow")
            if shadow and shadow.get("has_shadow"):
                visual_effects["has_shadows"] = True
                visual_effects["shadow_count"] += 1
                shadow_params.append({
                    "blur_pt": shadow.get("blur_pt", 4),
                    "offset_pt": shadow.get("offset_pt", 2),
                    "alpha_pct": shadow.get("alpha_pct", 25),
                    "color": shadow.get("color", "000000"),
                })
            elif shape.get("has_shadow"):
                visual_effects["has_shadows"] = True
                visual_effects["shadow_count"] += 1
            # Corner radius
            if shape.get("corner_radius_pct"):
                visual_effects["has_rounded_corners"] = True
                corner_radii.append(shape["corner_radius_pct"])
            # Transparency
            if shape.get("opacity_pct") and shape["opacity_pct"] < 100:
                visual_effects["has_transparency"] = True
            # Gradients with angle
            fill = shape.get("fill", {})
            if isinstance(fill, dict) and fill.get("gradient_stops"):
                visual_effects["has_gradients"] = True
                visual_effects["gradient_count"] += 1
                if fill.get("gradient_angle_deg") is not None:
                    gradient_angles.append(fill["gradient_angle_deg"])
            st = shape.get("shape_type", "")
            if st and st not in visual_effects["shape_variety"]:
                visual_effects["shape_variety"].append(st)
            # Detect decorative: shapes with fill but no text/table/chart/image
            has_content = bool(shape.get("text_frame", {}).get("paragraphs"))
            has_media = shape.get("table") or shape.get("chart") or shape.get("image")
            if not has_content and not has_media and fill:
                pos = shape.get("position", {})
                visual_effects["decorative_elements"].append({
                    "shape_type": st,
                    "fill": fill,
                    "position_pct": {
                        "left": pos.get("left_pct"),
                        "top": pos.get("top_pct"),
                        "width": pos.get("width_pct"),
                        "height": pos.get("height_pct"),
                    },
                })
    # Aggregate shadow style (median of observed parameters)
    if shadow_params:
        visual_effects["shadow_style"] = {
            "blur_pt": round(sorted(p["blur_pt"] for p in shadow_params)[len(shadow_params)//2], 1),
            "offset_pt": round(sorted(p["offset_pt"] for p in shadow_params)[len(shadow_params)//2], 1),
            "alpha_pct": round(sorted(p["alpha_pct"] for p in shadow_params)[len(shadow_params)//2], 1),
            "color": shadow_params[0]["color"],  # Most common shadow color
        }
    if corner_radii:
        visual_effects["corner_radius_pct"] = round(sum(corner_radii) / len(corner_radii), 1)
    if gradient_angles:
        visual_effects["gradient_angles"] = list(set(gradient_angles))
    result["visual_effects"] = visual_effects

    # ── Design Philosophy Analysis ──
    # Analyze COMPOSITION: how elements relate to each other, not just what they are.
    # This is what separates "list of shapes" from "understanding the design."
    philosophy = {}

    # 1. Content density: how much text per slide on average?
    text_counts = []
    for slide in slides_data:
        char_count = 0
        for shape in slide.get("shapes", []):
            tf = shape.get("text_frame", {})
            for para in tf.get("paragraphs", []):
                for run in para.get("runs", []):
                    char_count += len(run.get("text", ""))
        text_counts.append(char_count)
    if text_counts:
        avg_chars = sum(text_counts) / len(text_counts)
        philosophy["content_density"] = (
            "minimal" if avg_chars < 80 else
            "moderate" if avg_chars < 200 else
            "dense" if avg_chars < 400 else
            "very_dense"
        )
        philosophy["avg_chars_per_slide"] = round(avg_chars)

    # 2. Whitespace ratio: how much of each slide is empty?
    coverage_pcts = []
    for slide in slides_data:
        total_area = 0
        for shape in slide.get("shapes", []):
            pos = shape.get("position", {})
            w = pos.get("width_pct", 0) or 0
            h = pos.get("height_pct", 0) or 0
            total_area += w * h / 100  # as percentage of slide
        coverage_pcts.append(min(total_area, 100))
    if coverage_pcts:
        avg_coverage = sum(coverage_pcts) / len(coverage_pcts)
        philosophy["whitespace_ratio"] = round(100 - avg_coverage)
        philosophy["whitespace_style"] = (
            "airy" if avg_coverage < 30 else
            "balanced" if avg_coverage < 55 else
            "compact" if avg_coverage < 75 else
            "dense"
        )

    # 3. Shape count per slide: simple vs complex layouts
    shape_counts = [slide.get("shape_count", 0) for slide in slides_data]
    if shape_counts:
        avg_shapes = sum(shape_counts) / len(shape_counts)
        philosophy["avg_shapes_per_slide"] = round(avg_shapes, 1)
        philosophy["layout_complexity"] = (
            "minimal" if avg_shapes < 4 else
            "clean" if avg_shapes < 7 else
            "moderate" if avg_shapes < 12 else
            "complex"
        )

    # 4. Text alignment preference
    align_counts = Counter()
    for slide in slides_data:
        for shape in slide.get("shapes", []):
            tf = shape.get("text_frame", {})
            for para in tf.get("paragraphs", []):
                a = para.get("alignment", "left")
                align_counts[a] += 1
    if align_counts:
        philosophy["dominant_alignment"] = align_counts.most_common(1)[0][0]
        philosophy["alignment_distribution"] = dict(align_counts.most_common())

    # 5. Color usage restraint
    unique_colors = set(c for c in all_colors if c and not c.startswith("theme:"))
    philosophy["color_palette_size"] = len(unique_colors)
    philosophy["color_restraint"] = (
        "monochrome" if len(unique_colors) <= 3 else
        "restrained" if len(unique_colors) <= 6 else
        "moderate" if len(unique_colors) <= 12 else
        "colorful"
    )

    # 6. Visual effects restraint
    total_shapes = sum(shape_counts) if shape_counts else 1
    philosophy["effects_usage"] = {
        "shadows": "none" if not visual_effects["has_shadows"] else
                   "sparse" if visual_effects["shadow_count"] / max(total_shapes, 1) < 0.1 else
                   "moderate" if visual_effects["shadow_count"] / max(total_shapes, 1) < 0.3 else
                   "heavy",
        "gradients": "none" if not visual_effects["has_gradients"] else
                     "sparse" if visual_effects["gradient_count"] / max(total_shapes, 1) < 0.1 else
                     "moderate" if visual_effects["gradient_count"] / max(total_shapes, 1) < 0.3 else
                     "heavy",
    }

    # 7. Layout variety: does the deck use the same layout or mix it up?
    layout_names = [s.get("layout_name") for s in slides_data if s.get("layout_name")]
    unique_layouts = len(set(layout_names))
    philosophy["layout_variety"] = (
        "uniform" if unique_layouts <= 2 else
        "moderate" if unique_layouts <= 4 else
        "varied"
    )

    # 8. Title style: are titles big/centered or small/left-aligned?
    title_aligns = []
    title_sizes = []
    for slide in slides_data:
        shapes = slide.get("shapes", [])
        if not shapes:
            continue
        # Find likely title (largest text, topmost)
        texts = []
        for s in shapes:
            tf = s.get("text_frame", {})
            for p in tf.get("paragraphs", []):
                for r in p.get("runs", []):
                    sz = r.get("size_pt", 0)
                    if sz > 0:
                        texts.append((sz, p.get("alignment", "left"), s.get("position", {})))
        if texts:
            texts.sort(key=lambda x: x[0], reverse=True)
            title_sizes.append(texts[0][0])
            title_aligns.append(texts[0][1])
    if title_aligns:
        philosophy["title_alignment"] = Counter(title_aligns).most_common(1)[0][0]
    if title_sizes:
        philosophy["avg_title_size_pt"] = round(sum(title_sizes) / len(title_sizes), 1)

    result["design_philosophy"] = philosophy

    return result


def to_markdown(data):
    """Convert extracted data to a readable markdown design language doc."""
    lines = []
    lines.append(f"# Design Language: {data['source_file']}")
    lines.append("")

    # Dimensions
    dim = data["slide_dimensions"]
    lines.append(f"## Slide Canvas")
    lines.append(f"- **Dimensions**: {dim['width_in']}\" × {dim['height_in']}\" (aspect {dim['aspect_ratio']})")
    lines.append(f"- **Total slides**: {data['total_slides']}")
    lines.append("")

    # Theme
    theme = data.get("theme", {})
    if theme.get("color_scheme"):
        lines.append("## Theme Color Scheme")
        for name, color in theme["color_scheme"].items():
            lines.append(f"- **{name}**: `{color}`")
        lines.append("")

    # Design system
    ds = data.get("design_system", {})

    if ds.get("color_palette"):
        lines.append("## Color Palette (by frequency)")
        for c in ds["color_palette"]:
            lines.append(f"- `{c['hex']}` — used {c['uses']}×")
        lines.append("")

    if ds.get("typography"):
        typo = ds["typography"]
        if typo.get("fonts_by_frequency"):
            lines.append("## Fonts")
            for f in typo["fonts_by_frequency"]:
                lines.append(f"- **{f['font']}** — used {f['uses']}×")
            lines.append("")
        if typo.get("sizes_by_frequency"):
            lines.append("## Text Sizes")
            for s in typo["sizes_by_frequency"]:
                lines.append(f"- **{s['size_pt']}pt** — used {s['uses']}×")
            lines.append("")

    if ds.get("text_hierarchy"):
        lines.append("## Text Hierarchy (inferred)")
        for h in ds["text_hierarchy"]:
            lines.append(f"- **{h['likely_role']}**: {h['size_pt']}pt ({h['uses']}× occurrences)")
        lines.append("")

    # Layout usage
    if data.get("layout_usage"):
        lines.append("## Layout Usage")
        for l in data["layout_usage"]:
            lines.append(f"- **{l['layout']}**: {l['slides_using']} slides")
        lines.append("")

    # Per-slide blueprint
    lines.append("## Slide-by-Slide Blueprint")
    lines.append("")
    for slide in data.get("slides", []):
        lines.append(f"### Slide {slide['slide_number']} — Layout: {slide.get('layout_name', 'unknown')}")
        if slide.get("background"):
            bg = slide["background"]
            lines.append(f"- **Background**: {bg.get('type', 'none')}, color: {bg.get('color', 'n/a')}")

        for shape in slide.get("shapes", []):
            pos = shape.get("position", {})
            pos_str = f"({pos.get('left_pct', '?')}%, {pos.get('top_pct', '?')}%) → {pos.get('width_pct', '?')}% × {pos.get('height_pct', '?')}%"

            shape_type = shape.get("shape_type", "")
            name = shape.get("name", "")

            # Placeholder
            ph = shape.get("placeholder")
            if ph:
                lines.append(f"- **[{ph['type']}]** `{name}` at {pos_str}")
            else:
                lines.append(f"- **{shape_type}** `{name}` at {pos_str}")

            # Fill
            fill = shape.get("fill")
            if fill and fill.get("color"):
                lines.append(f"  - Fill: `{fill['color']}`")

            # Text
            tf = shape.get("text_frame", {})
            for para in tf.get("paragraphs", []):
                for run in para.get("runs", []):
                    text_preview = run["text"][:60] + ("..." if len(run["text"]) > 60 else "")
                    font_str = run.get("font", "?")
                    size_str = run.get("size_pt", "?")
                    style = ""
                    if run.get("bold"):
                        style += "B"
                    if run.get("italic"):
                        style += "I"
                    lines.append(f"  - Text: \"{text_preview}\" — {font_str} {size_str}pt {style} {run.get('color', '')}".rstrip())

            # Table
            if shape.get("table"):
                t = shape["table"]
                lines.append(f"  - Table: {t['rows']}×{t['cols']}")

            # Image
            if shape.get("image"):
                img = shape["image"]
                lines.append(f"  - Image: {img.get('content_type', 'unknown')}")

            # Chart
            if shape.get("chart"):
                ch = shape["chart"]
                lines.append(f"  - Chart: {ch.get('chart_type', 'unknown')}, {ch.get('series_count', 0)} series")

        lines.append("")

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description="Extract design language from a PowerPoint file",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python ppt_design_extractor.py presentation.pptx
  python ppt_design_extractor.py deck.pptx -o design.json
  python ppt_design_extractor.py deck.pptx --format markdown -o design.md
  python ppt_design_extractor.py deck.pptx --format both -o design
        """
    )
    parser.add_argument("input", help="Path to .pptx file")
    parser.add_argument("-o", "--output", help="Output file path (default: stdout for json, auto-named for markdown)")
    parser.add_argument("--format", choices=["json", "markdown", "both"], default="both",
                        help="Output format (default: both)")

    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    print(f"Extracting design language from: {args.input}", file=sys.stderr)
    data = extract_ppt(args.input)

    if args.format in ("json", "both"):
        json_out = args.output if args.format == "json" else (f"{args.output}.json" if args.output else None)
        json_str = json.dumps(data, indent=2, ensure_ascii=False)
        if json_out:
            Path(json_out).write_text(json_str, encoding="utf-8")
            print(f"JSON written to: {json_out}", file=sys.stderr)
        else:
            print(json_str)

    if args.format in ("markdown", "both"):
        md_out = args.output if args.format == "markdown" else (f"{args.output}.md" if args.output else None)
        md_str = to_markdown(data)
        if md_out:
            Path(md_out).write_text(md_str, encoding="utf-8")
            print(f"Markdown written to: {md_out}", file=sys.stderr)
        else:
            print(md_str)

    # Summary
    ds = data.get("design_system", {})
    colors = len(ds.get("color_palette", []))
    fonts = len(ds.get("typography", {}).get("fonts_by_frequency", []))
    print(f"\nExtracted: {data['total_slides']} slides, {colors} colors, {fonts} fonts", file=sys.stderr)


if __name__ == "__main__":
    main()
